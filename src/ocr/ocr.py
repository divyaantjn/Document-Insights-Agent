from __future__ import annotations

import asyncio
import base64
import gc
import io
import os
import uuid
import zipfile
from concurrent.futures import ThreadPoolExecutor
from contextlib import suppress
from io import BytesIO
from typing import Any, Dict, List, Optional, Tuple

import aiofiles
import pandas as pd
import pdfplumber
import yaml
from docx import Document
from docx.table import Table as DocxTable
from langchain_core.documents import Document as LangchainDocument
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader
from pypdf.errors import PdfReadError

from src.llm.litellm_client import LitellmClient
from src.ocr.semantic_chunker import SemanticChunker
from src.utils.logger import setup_logger
from src.utils.reasoning_extractor import REASONING_SECTION_PROMPT

import io as _io

import atexit
import threading
import zipfile as _zipfile_module

# ── Module-level singletons ───────────────────────────────────────────────────
litellm_client = LitellmClient()

# Raised from 4 → 8 to reduce starvation under concurrent Lambda cold-starts.
# wait=False: on Lambda shutdown the executor will not block process exit
# waiting for in-flight threads.  cancel_futures=True drops any queued work.
_OCR_THREAD_EXECUTOR = ThreadPoolExecutor(
    max_workers=8, thread_name_prefix="ocr_worker"
)
atexit.register(lambda: _OCR_THREAD_EXECUTOR.shutdown(wait=False, cancel_futures=True))

# ── Validation limits ─────────────────────────────────────────────────────────
# Max uncompressed size of a ZIP-based office doc (DOCX/XLSX/PPTX) in bytes.
# Checked before any XML parsing to reject adversarially large files early.
_MAX_UNCOMPRESSED_BYTES: int = 50 * 1024 * 1024   # 50 MB

# Max total chars allowed in a document before we skip it as garbage.
# Real documents of 10 M+ chars with zero tables or images are synthetic.
_MAX_TOTAL_CHARS: int = 8_000_000   # 8 MB of text

# Timeouts (seconds) for blocking library open() calls run in the executor.
_DOCX_OPEN_TIMEOUT: int = 30
_PPTX_OPEN_TIMEOUT: int = 30
_EXCEL_OPEN_TIMEOUT: int = 30
_PDF_OPEN_TIMEOUT: int = 30      # hard cap on pdfplumber.open()
_PDF_IS_TEXT_TIMEOUT: int = 15   # hard cap on _pdf_is_text_bearing()

# ── Named constants (replaces magic numbers flagged by SonarLint) ─────────────
MILVUS_MAX_TEXT_LENGTH: int = 65_535
"""Hard text-length cap enforced by the Milvus VARCHAR field."""

_MAX_ELEMENT_CHARS: int = 50_000
"""
Maximum characters allowed in a single text element before it is pre-split
by RecursiveCharacterTextSplitter *before* SemanticChunker is invoked.
Prevents the 10 MB DOCX OOM scenario.
"""

_DOCX_FLUSH_CHARS: int = 5_000
"""
DOCX text buffer is flushed into elements every this many characters,
regardless of whether a table boundary is encountered.  Prevents unbounded
growth of a single text element when a document has no tables.
"""

_PDF_TEXT_DETECT_MIN_CHARS: int = 20
"""
Minimum characters returned by pypdf on page 0 to classify a PDF as
text-bearing.  Below this threshold the PDF is treated as image-only and the
60-second pdfplumber text-extraction call is skipped entirely.
"""


# ── Processing limits ─────────────────────────────────────────────────────────

class ProcessingLimits:
    """Centralised hard limits for document processing."""

    MAX_CONCURRENT_IMAGES: int = 3
    MAX_ZIP_DEPTH: int = 3
    MAX_ZIP_FILES: int = 100
    MAX_ZIP_EXTRACTED_SIZE: int = 100 * 1024 * 1024   # 100 MB
    MAX_PDF_IMAGES: int = 30
    MAX_PPTX_IMAGES: int = 30
    MAX_DOCX_IMAGES: int = 30
    MAX_PDF_IMAGES_PER_PAGE: int = 5
    IMAGE_ANALYSIS_TIMEOUT: int = 30
    IMAGE_ANALYSIS_TIMEOUT_PER_BATCH_CAP: int = 300
    MAX_PDF_PAGES: int = 500
    PDF_IMAGE_DPI: int = 72


# ── Table utilities ───────────────────────────────────────────────────────────

class TableProcessor:
    """Utility helpers for table extraction and markdown conversion."""

    @staticmethod
    def table_to_markdown(table_data: List[List[str]]) -> str:
        """Convert table_data (list of rows) to a GitHub-flavoured markdown table."""
        if not table_data:
            return ""

        header = table_data[0]
        header_str = "| " + " | ".join(str(cell) for cell in header) + " |"
        separator = "| " + " | ".join(["---"] * len(header)) + " |"

        rows: List[str] = []
        for row in table_data[1:]:
            padded = list(row) + [""] * (len(header) - len(row))
            rows.append("| " + " | ".join(str(c) for c in padded[: len(header)]) + " |")

        return "\n".join([header_str, separator] + rows)

    @staticmethod
    def chunk_large_table(
        table_data: List[List[str]], rows_per_chunk: int = 20
    ) -> List[str]:
        """Split large tables into smaller markdown chunks preserving the header row."""
        if not table_data or len(table_data) <= 1:
            return []

        header = table_data[0]
        data_rows = table_data[1:]

        if len(data_rows) <= rows_per_chunk:
            return [TableProcessor.table_to_markdown(table_data)]

        chunks: List[str] = []
        for start in range(0, len(data_rows), rows_per_chunk):
            sub = [header] + data_rows[start: start + rows_per_chunk]
            chunks.append(TableProcessor.table_to_markdown(sub))
        return chunks

    @staticmethod
    def extract_docx_table(table: DocxTable) -> List[List[str]]:
        """Extract data from a python-docx Table object."""
        return [[cell.text.strip() for cell in row.cells] for row in table.rows]

    @staticmethod
    def extract_pptx_table(table: Any) -> List[List[str]]:
        """Extract data from a python-pptx table object."""
        return [[cell.text.strip() for cell in row.cells] for row in table.rows]


# ── Helper: open pypdf reader ─────────────────────────────────────────────────

def _open_pypdf_reader(pdf_bytes: bytes) -> Optional[PdfReader]:
    """
    Open a pypdf PdfReader from raw bytes.
    Pure-Python, no native rendering, no poppler dependency.
    Returns None on any failure so callers can degrade gracefully.
    """
    try:
        return PdfReader(_io.BytesIO(pdf_bytes), strict=False)
    except (PdfReadError, Exception):  # noqa: BLE001
        return None


def _validate_zip_based_file(file_bytes: bytes, filename: str) -> Optional[str]:
    """
    Pre-flight validation for ZIP-based office formats (DOCX, XLSX, PPTX).

    Opens the ZIP central directory only (no XML parsing) and checks the
    total uncompressed size.  Returns an error string if the file should be
    rejected, or None if it passes.

    This catches adversarially large files (e.g. large.docx with 10 MB of
    random text) before any library parsing begins, at microsecond cost.
    """
    try:
        with _zipfile_module.ZipFile(_io.BytesIO(file_bytes), "r") as zf:
            total_uncompressed = sum(info.file_size for info in zf.infolist())
        if total_uncompressed > _MAX_UNCOMPRESSED_BYTES:
            return (
                f"{filename} is too large to process: uncompressed size "
                f"{total_uncompressed // (1024*1024)} MB exceeds the "
                f"{_MAX_UNCOMPRESSED_BYTES // (1024*1024)} MB limit. "
                f"Please split or reduce the file size."
            )
        return None
    except _zipfile_module.BadZipFile:
        return f"{filename} appears to be corrupt or is not a valid Office document."
    except Exception:  # noqa: BLE001
        return None   # Unknown error — let the parser decide


def _pdf_is_text_bearing(pypdf_reader: PdfReader) -> bool:
    """
    Return True when the PDF contains a selectable text layer.

    Uses per-page daemon threads with a per-page time budget so a single
    stuck page cannot block indefinitely.  The total budget is
    _PDF_IS_TEXT_TIMEOUT seconds split evenly across up to 3 pages.

    Daemon threads are used intentionally: if they are still alive when
    the function returns (timeout hit) they will be orphaned but will not
    block Python process shutdown because they are daemon threads.
    """
    pages_to_check = min(3, len(pypdf_reader.pages))
    collected_chars = 0
    per_page_budget = _PDF_IS_TEXT_TIMEOUT / max(pages_to_check, 1)

    result_holder: List[str] = []

    def _extract(page_idx: int) -> None:
        try:
            text = pypdf_reader.pages[page_idx].extract_text() or ""
            result_holder.append(text)
        except Exception:  # noqa: BLE001
            result_holder.append("")

    for idx in range(pages_to_check):
        result_holder.clear()
        t = threading.Thread(target=_extract, args=(idx,), daemon=True)
        t.start()
        t.join(timeout=per_page_budget)
        if t.is_alive():
            # Thread still blocked — PDF content stream is unreadable in budget.
            # Treat as image-only to avoid the 60s pdfplumber stall.
            return False
        if result_holder:
            collected_chars += len(result_holder[0].strip())

    return collected_chars >= _PDF_TEXT_DETECT_MIN_CHARS


# ── Main OCR class ────────────────────────────────────────────────────────────

class Ocr:
    """
    Multi-format document text extractor.

    Supports: DOCX, PPTX, PDF, XLSX, XLS, CSV, plain-text, images, ZIP.
    """

    def __init__(self, config_path: str = "config/ocr_config.yaml") -> None:
        self.logger = setup_logger()
        self.config = self._load_config(config_path)

        ocr_cfg: Dict[str, Any] = self.config.get("ocr", {})
        self.image_analysis_config: Dict[str, Any] = ocr_cfg.get("image_analysis", {})
        self.document_extraction_config: Dict[str, Any] = ocr_cfg.get("document_extraction", {})
        self.output_formatting_config: Dict[str, Any] = ocr_cfg.get("output_formatting", {})
        self.error_messages: Dict[str, Any] = ocr_cfg.get("error_messages", {})
        self.success_messages: Dict[str, Any] = ocr_cfg.get("success_messages", {})
        self.info_messages: Dict[str, Any] = ocr_cfg.get("info_messages", {})
        self.warning_messages: Dict[str, Any] = ocr_cfg.get("warning_messages", {})
        self.text_extraction_config: Dict[str, Any] = ocr_cfg.get("text_extraction", {})
        self.pdf_processing_config: Dict[str, Any] = ocr_cfg.get("pdf_processing", {})
        self.chunking_config: Dict[str, Any] = ocr_cfg.get("chunking", {})

        # Log-message templates (fallbacks when YAML keys are absent)
        self._MSG_PROCESSING_IMAGES = "Processing {count} images for comprehensive understanding"
        self._MSG_IMAGE_ANALYSIS_FAILED = "Image {idx} analysis failed: {error}"
        self._MSG_IMAGE_EXTRACTION_FAILED = "Failed to extract image {idx}: {error}"

        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunking_config.get("chunk_size", 1000),
            chunk_overlap=self.chunking_config.get("chunk_overlap", 200),
        )

        semantic_cfg: Dict[str, Any] = self.chunking_config.get("semantic_chunking", {})
        self.semantic_chunking_enabled: bool = semantic_cfg.get("enabled", True)
        self.semantic_chunker = SemanticChunker(
            embedding_model=semantic_cfg.get("embedding_model", "gemini-embedding-001"),
            task_type=semantic_cfg.get("task_type", "SEMANTIC_SIMILARITY"),
            output_dimensionality=semantic_cfg.get("output_dimensionality", 768),
            batch_size=semantic_cfg.get("batch_size", 32),
            percentile_threshold=semantic_cfg.get("percentile_threshold", 25.0),
            max_chunk_size=semantic_cfg.get("max_chunk_size", 2000),
            min_chunk_size=semantic_cfg.get("min_chunk_size", 100),
        )

        self._pdf_page_batch_size: int = self.chunking_config.get("pdf_page_batch_size", 10)
        self._image_semaphore = asyncio.Semaphore(ProcessingLimits.MAX_CONCURRENT_IMAGES)

    # ── Config ────────────────────────────────────────────────────────────────

    def _load_config(self, config_path: str) -> Dict[str, Any]:
        """Load configuration from YAML file, returning safe defaults on failure."""
        if not os.path.exists(config_path):
            self.logger.warning(
                "Configuration file not found at %s, using defaults", config_path
            )
            return {"ocr": {}}
        try:
            with open(config_path, "r", encoding="utf-8") as fh:
                config = yaml.safe_load(fh)
            self.logger.info("Loaded configuration from %s", config_path)
            return config or {"ocr": {}}
        except Exception:  # noqa: BLE001
            self.logger.error("Failed to load configuration from %s, using defaults", config_path)
            return {"ocr": {}}

    # ── Image relevance filter ────────────────────────────────────────────────

    def _is_relevant_image(
        self,
        image_bytes: bytes,
        width: Optional[int] = None,
        height: Optional[int] = None,
    ) -> bool:
        """Return True when an image is likely meaningful (chart, diagram, photo)."""
        try:
            min_width = self.image_analysis_config.get("min_image_width", 100)
            min_height = self.image_analysis_config.get("min_image_height", 100)
            min_file_size = self.image_analysis_config.get("min_image_file_size", 5000)
            max_aspect_ratio = self.image_analysis_config.get("max_aspect_ratio", 20)

            if len(image_bytes) < min_file_size:
                return False
            if width is not None and height is not None:
                if width < min_width or height < min_height:
                    return False
                if min(width, height) > 0:
                    aspect = max(width, height) / min(width, height)
                    if aspect > max_aspect_ratio:
                        return False
            return True
        except Exception:  # noqa: BLE001
            self.logger.warning("Image filter error — including image by default")
            return True

    # ── Pre-flight size splitter ──────────────────────────────────────────────

    def _presplit_large_text(self, text: str) -> List[str]:
        """
        Split a text element that exceeds _MAX_ELEMENT_CHARS into smaller
        pieces using the RecursiveCharacterTextSplitter *before* any
        semantic chunking attempt.

        This is the primary guard against the OOM scenario: a 10 MB DOCX
        body is split into ~200 chunks of ≤5 000 chars each before any
        Gemini API call is made.
        """
        if len(text) <= _MAX_ELEMENT_CHARS:
            return [text]
        self.logger.warning(
            "Text element %d chars exceeds pre-split limit %d — "
            "applying RecursiveCharacterTextSplitter before semantic chunking.",
            len(text),
            _MAX_ELEMENT_CHARS,
        )
        docs = self.text_splitter.create_documents([text])
        return [d.page_content for d in docs]

    # ── Chunking helpers ──────────────────────────────────────────────────────

    async def _chunk_elements(
        self,
        elements: List[Dict[str, Any]],
        source: str,
        page: int = 1,
        use_semantic: bool = True,
    ) -> List[LangchainDocument]:
        """Chunk extracted elements into LangchainDocuments."""
        chunks: List[LangchainDocument] = []
        rows_per_chunk = self.chunking_config.get("table_rows_per_chunk", 20)
        text_tasks = []

        for element in elements:
            metadata_base = {"source": source, "page": page, **element.get("metadata", {})}
            self._process_single_element_for_chunking(
                element, metadata_base, rows_per_chunk, use_semantic, chunks, text_tasks
            )

        if text_tasks:
            text_results = await self._gather_text_chunks(text_tasks)
            for result in text_results:
                if isinstance(result, Exception):
                    self.logger.error("Text chunking task failed: %s", result)
                    continue
                chunks.extend(result)

        return chunks

    def _process_single_element_for_chunking(
        self,
        element: Dict[str, Any],
        metadata_base: Dict[str, Any],
        rows_per_chunk: int,
        use_semantic: bool,
        chunks: List[LangchainDocument],
        text_tasks: List,
    ) -> None:
        """
        Dispatch a single element to the correct chunking strategy.

        Text elements that exceed _MAX_ELEMENT_CHARS are pre-split first so
        the semantic chunker never receives a document-scale string.
        """
        element_type = element.get("type")
        content = element.get("content")

        if element_type == "table":
            table_chunks = TableProcessor.chunk_large_table(content, rows_per_chunk)
            for chunk_idx, table_markdown in enumerate(table_chunks):
                chunks.append(
                    LangchainDocument(
                        page_content=table_markdown,
                        metadata={
                            **metadata_base,
                            "type": "table",
                            "table_chunk_index": chunk_idx if len(table_chunks) > 1 else None,
                        },
                    )
                )

        elif element_type == "text":
            # ── Pre-flight size guard ──────────────────────────────────
            sub_texts = self._presplit_large_text(content)
            for sub_text in sub_texts:
                text_tasks.append(
                    self._chunk_text_element(sub_text, metadata_base, use_semantic)
                )

        elif element_type == "image":
            chunks.append(
                LangchainDocument(
                    page_content=content,
                    metadata={**metadata_base, "type": "image"},
                )
            )

    async def _gather_text_chunks(self, text_tasks: List) -> List:
        """Gather text chunks with bounded concurrency."""
        sem = asyncio.Semaphore(
            self.chunking_config.get("max_concurrent_text_chunks", 10)
        )

        async def _limited(task: Any) -> Any:
            async with sem:
                return await task

        return await asyncio.gather(
            *[_limited(t) for t in text_tasks], return_exceptions=True
        )

    async def _chunk_text_element(
        self,
        content: str,
        metadata_base: Dict[str, Any],
        use_semantic: bool,
    ) -> List[LangchainDocument]:
        """Chunk a single text element, falling back to recursive splitter on failure."""
        if use_semantic and self.semantic_chunking_enabled:
            try:
                timeout_val = self.chunking_config.get("semantic_chunking_timeout", 30)
                async with asyncio.timeout(timeout_val):
                    semantic_texts = await self.semantic_chunker.chunk(content)
                if semantic_texts:
                    return [
                        LangchainDocument(
                            page_content=chunk_text,
                            metadata={**metadata_base, "type": "text"},
                        )
                        for chunk_text in semantic_texts
                    ]
            except asyncio.TimeoutError:
                self.logger.warning(
                    "Semantic chunking timed out for %d-char element; "
                    "using recursive splitter.",
                    len(content),
                )
            except Exception:  # noqa: BLE001
                self.logger.warning(
                    "Semantic chunking failed; using recursive splitter.",
                    exc_info=True,
                )

        split_docs = self.text_splitter.create_documents(
            [content], metadatas=[{**metadata_base, "type": "text"}]
        )
        return list(split_docs)

    def _run_async_chunking(
        self,
        elements: List[Dict],
        source: str,
        use_semantic: bool = True,
    ) -> List[LangchainDocument]:
        """
        Run _chunk_elements safely from a synchronous context.

        Uses a fresh background thread when an event loop is already running
        (e.g. inside FastAPI / Lambda) to avoid RuntimeError: nested loop.
        """
        loop = asyncio.get_event_loop()
        if loop.is_running():
            with ThreadPoolExecutor(max_workers=1) as pool:
                future = pool.submit(
                    asyncio.run,
                    self._chunk_elements(elements, source=source, use_semantic=use_semantic),
                )
                return future.result()
        return loop.run_until_complete(
            self._chunk_elements(elements, source=source, use_semantic=use_semantic)
        )

    # ── Image analysis formatting ─────────────────────────────────────────────

    def _format_image_analysis(
        self, result: Dict[str, Any], image_name: str
    ) -> str:
        """Format image analysis result dict into readable annotated text."""
        try:
            header_fmt = self.output_formatting_config.get(
                "image_analysis_header_format",
                "Image Content Analysis - {image_name}",
            )
            output_parts: List[str] = [
                f"\n[{header_fmt.format(image_name=image_name)}]:\n"
            ]

            section_keys = (
                ("ocr_text",
                 self.output_formatting_config.get("text_header", "Text Content")),
                ("visual_description",
                 self.output_formatting_config.get("visual_description_header", "Visual Description")),
                ("key_insights",
                 self.output_formatting_config.get("key_insights_header", "Key Insights")),
                ("data_points",
                 self.output_formatting_config.get("data_points_header", "Data Points")),
            )

            has_content = False
            for key, header in section_keys:
                if result.get(key):
                    output_parts.append(f"\n{header}:\n{result[key]}\n")
                    has_content = True

            return "".join(output_parts) if has_content else ""
        except Exception:  # noqa: BLE001
            self.logger.error(
                "Error formatting image analysis for %s", image_name, exc_info=True
            )
            return ""

    # ── DOCX extraction ───────────────────────────────────────────────────────

    async def extract_text_from_docx(
        self,
        docx_bytes: bytes,
        team_config: Dict[str, Any],
        auth_token: str,
    ) -> Dict[str, Any]:
        """
        Extract text from DOCX with table preservation.

        Key changes vs prior version:
        - Text buffer is flushed every _DOCX_FLUSH_CHARS characters,
          preventing a single unbounded text element (root cause of OOM).
        - Images are guarded by MAX_DOCX_IMAGES + per-image timeout.
        """
        doc = None
        elements: List[Dict[str, Any]] = []
        image_tasks: List = []
        image_metadata: List = []

        # ── Pre-flight validation ──────────────────────────────────────────
        validation_error = _validate_zip_based_file(docx_bytes, "DOCX")
        if validation_error:
            self.logger.warning("DOCX rejected by pre-flight check: %s", validation_error)
            return {"error": validation_error}

        try:
            loop = asyncio.get_event_loop()
            try:
                doc = await asyncio.wait_for(
                    loop.run_in_executor(
                        _OCR_THREAD_EXECUTOR,
                        lambda b=docx_bytes: Document(io.BytesIO(b)),
                    ),
                    timeout=_DOCX_OPEN_TIMEOUT,
                )
            except asyncio.TimeoutError:
                self.logger.error(
                    "DOCX open timed out after %ds — file may be corrupt",
                    _DOCX_OPEN_TIMEOUT,
                )
                return {"error": "DOCX file could not be opened within the time limit. The file may be corrupt."}
            docx_bytes = b""   # release raw bytes; do not del (may raise in finally)

            char_count = 0
            current_page = 1
            chars_per_page = 3000
            text_buffer: List[str] = []

            for block in doc.element.body:
                if block.tag.endswith("tbl"):
                    char_count, current_page = self._flush_text_buffer(
                        text_buffer, elements, current_page
                    )
                    text_buffer = []
                    self._process_docx_table_block(block, doc, elements, current_page)

                elif block.tag.endswith("p"):
                    char_count, current_page = self._process_docx_paragraph(
                        block, doc, text_buffer, char_count, chars_per_page
                    )
                    # ── Streaming flush: prevent unbounded text elements ────
                    if char_count >= _DOCX_FLUSH_CHARS:
                        char_count, current_page = self._flush_text_buffer(
                            text_buffer, elements, current_page
                        )
                        text_buffer = []

            if text_buffer:
                elements.append({
                    "type": "text",
                    "content": "\n".join(text_buffer),
                    "metadata": {"page_number": current_page},
                })
            text_buffer = []

            image_tasks, image_metadata = self._collect_docx_images(
                doc, elements, char_count, chars_per_page, team_config, auth_token
            )

            if image_tasks:
                batch_timeout = self._calc_batch_image_timeout(len(image_tasks))
                try:
                    async with asyncio.timeout(batch_timeout):
                        await self._process_image_results(
                            image_tasks, image_metadata, elements, "docx"
                        )
                except asyncio.TimeoutError:
                    self.logger.error(
                        "DOCX image batch timed out after %ds — skipping %d image(s)",
                        batch_timeout,
                        len(image_tasks),
                    )

            if not elements:
                error_msg = self.error_messages.get(
                    "docx_no_content", "DOCX contains no extractable content"
                )
                self.logger.error("%s", error_msg)
                return {"status_code": 400, "message": error_msg, "data": None}

            chunks = await self._chunk_elements(
                elements, source="docx_file", use_semantic=False
            )
            # use_semantic=False: DOCX is already paragraph-structured so
            # RecursiveCharacterTextSplitter is sufficient and avoids hundreds
            # of Gemini embedding API calls that caused 3-4 min extraction times.
            success_msg = self.success_messages.get(
                "docx_extracted", "Successfully extracted text from DOCX file"
            )
            self.logger.info("%s — %d chunks", success_msg, len(chunks))
            return {
                "chunks": [
                    {"text": c.page_content, **c.metadata} for c in chunks
                ]
            }

        except Exception:  # noqa: BLE001
            self.logger.error("Failed to extract text from DOCX", exc_info=True)
            return {"error": "Failed to extract text from DOCX"}

        finally:
            if doc is not None:
                with suppress(Exception):
                    del doc
            elements.clear()
            image_tasks.clear()
            image_metadata.clear()
            gc.collect()

    @staticmethod
    def _calc_batch_image_timeout(image_count: int) -> int:
        """Compute a bounded batch timeout based on image count."""
        return min(
            ProcessingLimits.IMAGE_ANALYSIS_TIMEOUT_PER_BATCH_CAP,
            max(60, image_count * ProcessingLimits.IMAGE_ANALYSIS_TIMEOUT),
        )

    def _flush_text_buffer(
        self,
        text_buffer: List[str],
        elements: List[Dict],
        current_page: int,
    ) -> Tuple[int, int]:
        """Flush the accumulated text buffer into elements and reset char count."""
        if text_buffer:
            elements.append({
                "type": "text",
                "content": "\n".join(text_buffer),
                "metadata": {"page_number": current_page},
            })
            text_buffer.clear()
        return 0, current_page

    def _process_docx_table_block(
        self,
        block: Any,
        doc: Document,
        elements: List[Dict],
        current_page: int,
    ) -> None:
        """Find the python-docx Table matching *block* and append its data."""
        for table in doc.tables:
            if table._element == block:  # noqa: SLF001 — no public API alternative
                table_data = TableProcessor.extract_docx_table(table)
                if table_data:
                    elements.append({
                        "type": "table",
                        "content": table_data,
                        "metadata": {"page_number": current_page},
                    })
                break

    def _process_docx_paragraph(
        self,
        block: Any,
        doc: Document,
        text_buffer: List[str],
        char_count: int,
        chars_per_page: int,
    ) -> Tuple[int, int]:
        """Append the paragraph text (if non-empty) to *text_buffer*."""
        for para in doc.paragraphs:
            if para._element == block:  # noqa: SLF001
                if para.text.strip():
                    text_buffer.append(para.text)
                    char_count += len(para.text)
                break
        current_page = (char_count // chars_per_page) + 1
        return char_count, current_page

    def _collect_docx_images(
        self,
        doc: Document,
        elements: List[Dict],
        char_count: int,
        chars_per_page: int,
        team_config: Dict,
        auth_token: str,
    ) -> Tuple[List, List]:
        """Collect image analysis tasks from DOCX inline shapes."""
        context_window = self.document_extraction_config.get("docx_context_window", 500)
        image_tasks: List = []
        image_metadata: List = []
        image_count = 0

        for idx, inline_shape in enumerate(doc.inline_shapes, 1):
            if image_count >= ProcessingLimits.MAX_DOCX_IMAGES:
                self.logger.warning(
                    "Reached maximum image limit (%d) for DOCX",
                    ProcessingLimits.MAX_DOCX_IMAGES,
                )
                break

            task, meta = self._extract_docx_inline_image(
                inline_shape, idx, doc, elements,
                char_count, chars_per_page, context_window,
                team_config, auth_token,
            )
            if task is not None:
                image_tasks.append(task)
                image_metadata.append(meta)
                image_count += 1

        return image_tasks, image_metadata

    def _extract_docx_inline_image(
        self,
        inline_shape: Any,
        idx: int,
        doc: Document,
        elements: List[Dict],
        char_count: int,
        chars_per_page: int,
        context_window: int,
        team_config: Dict,
        auth_token: str,
    ) -> Tuple[Optional[Any], Optional[Dict]]:
        """Extract a single inline image from a DOCX document."""
        try:
            blip = inline_shape._inline.graphic.graphicData.pic.blipFill.blip  # noqa: SLF001
            r_id = blip.embed
            image_part = doc.part.related_parts[r_id]
            image_bytes: bytes = image_part.image.blob
            width: int = image_part.image.px_width
            height: int = image_part.image.px_height

            if not self._is_relevant_image(image_bytes, width, height):
                self.logger.debug("Skipping decorative image %d in DOCX", idx)
                return None, None

            image_page = (char_count // chars_per_page) + 1
            meta = {
                "image_number": idx,
                "filename": image_part.filename,
                "content_type": image_part.content_type,
                "size": (width, height),
                "page_number": image_page,
            }

            recent_text = " ".join(
                e["content"] if e["type"] == "text" else ""
                for e in elements[-3:]
            )
            context = recent_text[-context_window:] if recent_text else ""

            task = self._analyze_image_with_semaphore(
                content=image_bytes,
                team_config=team_config,
                auth_token=auth_token,
                context=context,
            )
            image_bytes = b""  # release; let GC reclaim
            return task, meta

        except Exception:  # noqa: BLE001
            self.logger.warning(
                "Failed to extract image %d from DOCX", idx, exc_info=True
            )
            return None, None

    # ── Image result processing ───────────────────────────────────────────────

    async def _process_image_results(
        self,
        image_tasks: List,
        image_metadata: List,
        elements: List,
        file_type: str,
    ) -> None:
        """
        Process image analysis tasks and stream results into *elements*.

        Uses asyncio.as_completed so each result is consumed and released
        immediately — never holding all LLM responses in RAM simultaneously.
        Each task is also guarded by IMAGE_ANALYSIS_TIMEOUT.
        """
        if not image_tasks:
            return

        info_msg = self.info_messages.get(
            "processing_images", self._MSG_PROCESSING_IMAGES
        )
        self.logger.info("%s", info_msg.format(count=len(image_tasks)))

        per_img_timeout = ProcessingLimits.IMAGE_ANALYSIS_TIMEOUT

        async def _indexed_call(
            i: int, coro: Any
        ) -> Tuple[int, Any]:
            result = await self._guarded_image_call(coro, per_img_timeout)
            return i, result

        futures = [
            asyncio.ensure_future(_indexed_call(i, coro))
            for i, coro in enumerate(image_tasks)
        ]

        for future in asyncio.as_completed(futures):
            try:
                idx, result = await future
            except Exception:  # noqa: BLE001
                self.logger.error("Image analysis future raised unexpectedly", exc_info=True)
                continue

            if result is None or isinstance(result, Exception):
                self.logger.error(
                    "Image analysis failed for index %s: %s",
                    idx if not isinstance(result, Exception) else "?",
                    result,
                )
                continue

            if isinstance(result, dict):
                metadata = image_metadata[idx] if idx < len(image_metadata) else {}
                self._append_image_element(result, metadata, elements, file_type)

    @staticmethod
    async def _guarded_image_call(coro: Any, timeout_sec: float) -> Any:
        """Wrap a coroutine with a per-image timeout. Returns None on timeout."""
        try:
            async with asyncio.timeout(timeout_sec):
                return await coro
        except asyncio.TimeoutError:
            return None
        except Exception as exc:  # noqa: BLE001
            return exc

    def _append_image_element(
        self,
        result: Dict,
        metadata: Dict,
        elements: List[Dict],
        file_type: str,
    ) -> None:
        """Format and append a single image analysis result to elements."""
        if file_type == "docx":
            name = metadata.get("filename", f"Image {metadata.get('image_number', '?')}")
            page = metadata.get("page_number", 1)
            meta_out: Dict[str, Any] = {"page_number": page, "image_filename": name}
        elif file_type == "pptx":
            name = metadata.get("shape_name", f"Image {metadata.get('image_number', '?')}")
            page = metadata.get("slide_number", 1)
            meta_out = {"slide_number": page, "image_name": name}
        else:
            page = metadata.get("page_number", 1)
            img_num = metadata.get("image_number", 1)
            name = f"Page {page}, Image {img_num}"
            meta_out = {"page_number": page, "image_number": img_num}
            self.logger.info("[PDF Image Analysis] Appended: %s", name)

        analysis_output = self._format_image_analysis(result, name)
        if analysis_output:
            elements.append({
                "type": "image",
                "content": analysis_output,
                "metadata": meta_out,
            })

    # ── Excel / XLSX extraction ───────────────────────────────────────────────

    async def extract_text_from_excel(
        self, excel_bytes: bytes, filename: str
    ) -> Dict[str, Any]:
        """Extract text content from Excel files (.xlsx, .xls) with table preservation."""
        if not filename.endswith((".xlsx", ".xls")):
            return {"error": "Unsupported Excel file format"}

        # ── Pre-flight validation ──────────────────────────────────────────
        validation_error = _validate_zip_based_file(excel_bytes, filename)
        if validation_error:
            self.logger.warning("Excel rejected by pre-flight check: %s", validation_error)
            return {"error": validation_error}

        excel_file = None
        excel_buffer = None
        elements: List[Dict] = []

        try:
            excel_buffer = io.BytesIO(excel_bytes)
            excel_bytes = b""  # release raw bytes

            # Open with a timeout via executor to guard against corrupt files
            # that cause openpyxl/xlrd to hang during ZIP/XML traversal.
            loop = asyncio.get_event_loop()
            try:
                if loop.is_running():
                    import concurrent.futures
                    with concurrent.futures.ThreadPoolExecutor(max_workers=1) as _p:
                        excel_file = _p.submit(
                            pd.ExcelFile, excel_buffer
                        ).result(timeout=_EXCEL_OPEN_TIMEOUT)
                else:
                    excel_file = await asyncio.wait_for(
                        loop.run_in_executor(
                            _OCR_THREAD_EXECUTOR,
                            lambda buf=excel_buffer: pd.ExcelFile(buf),
                        ),
                        timeout=_EXCEL_OPEN_TIMEOUT,
                    )
            except (asyncio.TimeoutError, Exception) as exc:  # noqa: BLE001
                self.logger.error(
                    "Excel open timed out or failed for %s: %s", filename, exc
                )
                return {"error": f"Excel file could not be opened: the file may be corrupt or too large."}

            for sheet_name in excel_file.sheet_names:
                self._process_excel_sheet(excel_file, sheet_name, elements)

            if not elements:
                error_msg = self.error_messages.get(
                    "excel_no_content", "Excel does not contain data"
                )
                self.logger.error("%s", error_msg)
                return {"status_code": 400, "message": error_msg, "data": None}

            chunks = self._run_async_chunking(
                elements, source=filename, use_semantic=False
            )
            success_msg = self.success_messages.get(
                "excel_extracted",
                "Successfully extracted text from Excel file: {filename}",
            )
            self.logger.info("%s", success_msg.format(filename=filename))
            return {"chunks": [{"text": d.page_content, **d.metadata} for d in chunks]}

        except Exception:  # noqa: BLE001
            self.logger.error("Exception during Excel text extraction", exc_info=True)
            return {"error": "Failed to extract text from Excel file"}

        finally:
            if excel_file is not None:
                with suppress(Exception):
                    del excel_file
            if excel_buffer is not None:
                with suppress(Exception):
                    excel_buffer.close()
            elements.clear()

    def _process_excel_sheet(
        self,
        excel_file: pd.ExcelFile,
        sheet_name: str,
        elements: List[Dict],
    ) -> None:
        """Process a single Excel sheet, splitting large sheets into sub-table
        elements to prevent oversized Milvus/Kafka insert payloads.

        Each sub-element is capped at MAX_ROWS_PER_ELEMENT rows so that the
        resulting markdown chunk stays well under the 1 MB sub-batch limit in
        _insert_batch, even for wide sheets with many columns.
        """
        MAX_ROWS_PER_ELEMENT = 200  # ~200 rows keeps each table element under ~50 KB
        df = None
        try:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            table_data: List[List[str]] = [df.columns.tolist()] + df.values.tolist()
            table_data = [
                [str(cell) if pd.notna(cell) else "" for cell in row]
                for row in table_data
            ]
            if len(table_data) <= 1:
                return

            header = table_data[0]
            data_rows = table_data[1:]

            # Split large sheets into sub-elements, each with the header preserved
            for chunk_start in range(0, len(data_rows), MAX_ROWS_PER_ELEMENT):
                chunk_rows = data_rows[chunk_start: chunk_start + MAX_ROWS_PER_ELEMENT]
                elements.append({
                    "type": "table",
                    "content": [header] + chunk_rows,
                    "metadata": {
                        "sheet_name": sheet_name,
                        "row_offset": chunk_start,
                    },
                })
        except Exception:  # noqa: BLE001
            self.logger.error(
                "Failed to read sheet '%s'", sheet_name, exc_info=True
            )
        finally:
            if df is not None:
                with suppress(Exception):
                    del df

    async def extract_text_from_xlsx(self, excel_bytes: bytes) -> Dict[str, Any]:
        """Thin wrapper for standalone XLSX extraction (ZIP / direct calls)."""
        return await self.extract_text_from_excel(excel_bytes, "xlsx_file")

    # ── PPTX extraction ───────────────────────────────────────────────────────

    async def extract_text_from_ppt(
        self,
        ppt_bytes: bytes,
        team_config: Dict[str, Any],
        auth_token: str,
    ) -> Dict[str, Any]:
        """Extract text from PPTX with table and image preservation."""
        tmp_file_path: Optional[str] = None
        presentation = None
        elements: List[Dict] = []
        image_tasks: List = []
        image_metadata: List = []

        # ── Pre-flight validation ──────────────────────────────────────────
        validation_error = _validate_zip_based_file(ppt_bytes, "PPTX")
        if validation_error:
            self.logger.warning("PPTX rejected by pre-flight check: %s", validation_error)
            return {"error": validation_error}

        try:
            tmp_file_path = f"/tmp/pptx_{uuid.uuid4().hex}.pptx"
            async with aiofiles.open(tmp_file_path, "wb") as fh:
                await fh.write(ppt_bytes)
            ppt_bytes = b""  # release raw bytes

            loop = asyncio.get_event_loop()
            try:
                presentation = await asyncio.wait_for(
                    loop.run_in_executor(
                        _OCR_THREAD_EXECUTOR,
                        lambda p=tmp_file_path: Presentation(p),
                    ),
                    timeout=_PPTX_OPEN_TIMEOUT,
                )
            except asyncio.TimeoutError:
                self.logger.error(
                    "PPTX open timed out after %ds — file may be corrupt",
                    _PPTX_OPEN_TIMEOUT,
                )
                return {"error": "PPTX file could not be opened within the time limit. The file may be corrupt."}
            context_window = self.document_extraction_config.get("ppt_context_window", 300)
            image_count = 0

            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_tasks, slide_meta = self._process_pptx_slide(
                    slide, slide_num, elements, context_window,
                    image_count, team_config, auth_token,
                )
                image_tasks.extend(slide_tasks)
                image_metadata.extend(slide_meta)
                image_count += len(slide_tasks)

            if image_tasks:
                batch_timeout = self._calc_batch_image_timeout(len(image_tasks))
                try:
                    async with asyncio.timeout(batch_timeout):
                        await self._process_image_results(
                            image_tasks, image_metadata, elements, "pptx"
                        )
                except asyncio.TimeoutError:
                    self.logger.error(
                        "PPTX image batch timed out after %ds — skipping %d image(s)",
                        batch_timeout,
                        len(image_tasks),
                    )

            if not elements:
                error_msg = self.error_messages.get(
                    "ppt_no_content", "PPT contains no extractable content"
                )
                self.logger.error("%s", error_msg)
                return {"status_code": 400, "message": error_msg, "data": None}

            chunks = await self._chunk_elements(
                elements, source="ppt_file", use_semantic=False
            )
            # use_semantic=False: PPTX slides are short discrete text blocks,
            # paragraph-level splitting is sufficient and avoids Gemini API
            # calls per element that cause unnecessary latency.
            success_msg = self.success_messages.get(
                "ppt_extracted", "Successfully extracted text from PowerPoint file"
            )
            self.logger.info("%s — %d chunks", success_msg, len(chunks))
            return {
                "chunks": [
                    {"text": c.page_content, **c.metadata} for c in chunks
                ]
            }

        except Exception:  # noqa: BLE001
            self.logger.error("Failed to extract text from PowerPoint", exc_info=True)
            return {"error": "Failed to extract text from PowerPoint"}

        finally:
            if tmp_file_path and os.path.exists(tmp_file_path):
                with suppress(Exception):
                    os.unlink(tmp_file_path)
            if presentation is not None:
                with suppress(Exception):
                    del presentation
            elements.clear()
            image_tasks.clear()
            image_metadata.clear()
            gc.collect()

    def _process_pptx_slide(
        self,
        slide: Any,
        slide_num: int,
        elements: List[Dict],
        context_window: int,
        image_count: int,
        team_config: Dict,
        auth_token: str,
    ) -> Tuple[List, List]:
        """Process a single PPTX slide; return image tasks and metadata."""
        text_buffer: List[str] = []
        slide_image_tasks: List = []
        slide_image_metadata: List = []

        for shape in slide.shapes:
            if image_count + len(slide_image_tasks) >= ProcessingLimits.MAX_PPTX_IMAGES:
                self.logger.warning(
                    "Reached maximum image limit (%d) for PPTX",
                    ProcessingLimits.MAX_PPTX_IMAGES,
                )
                break
            self._process_pptx_shape(
                shape, slide_num, elements, text_buffer,
                slide_image_tasks, slide_image_metadata,
                context_window, team_config, auth_token,
            )

        if text_buffer:
            elements.append({
                "type": "text",
                "content": "\n".join(text_buffer),
                "metadata": {"slide_number": slide_num},
            })

        return slide_image_tasks, slide_image_metadata

    def _process_pptx_shape(
        self,
        shape: Any,
        slide_num: int,
        elements: List[Dict],
        text_buffer: List[str],
        image_tasks: List,
        image_metadata: List,
        context_window: int,
        team_config: Dict,
        auth_token: str,
    ) -> None:
        """Dispatch a PPTX shape to the correct handler."""
        if shape.has_table:
            self._flush_pptx_text_buffer(text_buffer, elements, slide_num)
            self._process_pptx_table_shape(shape, slide_num, elements)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            task, meta = self._extract_pptx_image(
                shape, slide_num, text_buffer, context_window, team_config, auth_token
            )
            if task is not None:
                image_tasks.append(task)
                image_metadata.append(meta)
        elif hasattr(shape, "text") and shape.text:
            text_buffer.append(shape.text)

    def _flush_pptx_text_buffer(
        self,
        text_buffer: List[str],
        elements: List[Dict],
        slide_num: int,
    ) -> None:
        """Flush the PPTX text buffer into an element and clear it."""
        if text_buffer:
            elements.append({
                "type": "text",
                "content": "\n".join(text_buffer),
                "metadata": {"slide_number": slide_num},
            })
            text_buffer.clear()

    def _process_pptx_table_shape(
        self,
        shape: Any,
        slide_num: int,
        elements: List[Dict],
    ) -> None:
        """Extract and append a table from a PPTX shape."""
        table_data = TableProcessor.extract_pptx_table(shape.table)
        if table_data:
            elements.append({
                "type": "table",
                "content": table_data,
                "metadata": {"slide_number": slide_num},
            })

    def _extract_pptx_image(
        self,
        shape: Any,
        slide_num: int,
        text_buffer: List[str],
        context_window: int,
        team_config: Dict,
        auth_token: str,
    ) -> Tuple[Optional[Any], Optional[Dict]]:
        """Extract a single image from a PPTX picture shape."""
        try:
            image = shape.image
            image_blob: bytes = image.blob
            width, height = image.size

            if not self._is_relevant_image(image_blob, width, height):
                self.logger.debug(
                    "Skipping decorative image '%s' on slide %d",
                    shape.name, slide_num,
                )
                return None, None

            meta = {
                "slide_number": slide_num,
                "shape_name": shape.name,
                "extension": image.ext,
                "size": image.size,
            }

            slide_context = " ".join(text_buffer)
            context = slide_context[-context_window:] if slide_context else ""

            task = self._analyze_image_with_semaphore(
                content=image_blob,
                team_config=team_config,
                auth_token=auth_token,
                context=context,
            )
            image_blob = b""  # release
            return task, meta

        except Exception:  # noqa: BLE001
            self.logger.warning(
                "Failed to extract image from slide %d", slide_num, exc_info=True
            )
            return None, None

    # ── PDF extraction ────────────────────────────────────────────────────────

    async def extract_text_from_pdf_file(
        self,
        content: bytes,
        team_config: Dict[str, Any],
        auth_token: str,
        chunk_callback: Any = None,
        max_pages: Optional[int] = None,
    ) -> Dict[str, Any]:
        """
        Extract text and images from a PDF file.

        New: scanned-image PDF detection via pypdf before pdfplumber.
        A PDF with no selectable text layer (< _PDF_TEXT_DETECT_MIN_CHARS on
        first 3 pages) skips the 60-second pdfplumber text-extraction path
        entirely, eliminating the timeout seen in the incident logs.

        STREAMING MODE (chunk_callback provided):
            Each page batch is chunked and passed to chunk_callback immediately.
        LEGACY MODE (chunk_callback=None):
            All chunks accumulated and returned.
        """
        pdf_doc = None
        pypdf_reader = None
        pdf_buffer = None
        all_chunks: List = []
        loop = asyncio.get_event_loop()

        try:
            pdf_doc, pdf_buffer, total_pages = await self._initialize_pdf_extraction(
                content, max_pages
            )
            if isinstance(pdf_doc, dict):   # error sentinel
                return pdf_doc  # type: ignore[return-value]

            # Open pypdf for image extraction + scanned-PDF classification
            try:
                pypdf_reader = await loop.run_in_executor(
                    _OCR_THREAD_EXECUTOR, _open_pypdf_reader, content
                )
            except Exception:  # noqa: BLE001
                self.logger.warning("pypdf open failed — image extraction disabled")
                pypdf_reader = None

            # ── Scanned-PDF detection with hard timeout ──────────────────────────────
            # _pdf_is_text_bearing uses per-page daemon threads capped at
            # _PDF_IS_TEXT_TIMEOUT total.  We also wrap the run_in_executor
            # call itself so a frozen executor thread never stalls the event loop.
            is_text_pdf = True
            if pypdf_reader is not None:
                try:
                    is_text_pdf = await asyncio.wait_for(
                        loop.run_in_executor(
                            _OCR_THREAD_EXECUTOR, _pdf_is_text_bearing, pypdf_reader
                        ),
                        timeout=_PDF_IS_TEXT_TIMEOUT + 5,
                    )
                except asyncio.TimeoutError:
                    self.logger.warning(
                        "PDF text-bearing check timed out after %ds — treating as image-only.",
                        _PDF_IS_TEXT_TIMEOUT + 5,
                    )
                    is_text_pdf = False
                if not is_text_pdf:
                    self.logger.warning(
                        "PDF classified as image-only (scanned) — "
                        "skipping pdfplumber text extraction to avoid timeout."
                    )

            # Both readers open — release original bytes
            content = b""

            context_window = self.document_extraction_config.get("pdf_context_window", 300)
            image_count = 0
            batch_size = self._pdf_page_batch_size

            _per_img_timeout = ProcessingLimits.IMAGE_ANALYSIS_TIMEOUT
            _batch_img_cap = min(
                batch_size * ProcessingLimits.MAX_PDF_IMAGES_PER_PAGE, 25
            )
            _batch_image_timeout = min(
                ProcessingLimits.IMAGE_ANALYSIS_TIMEOUT_PER_BATCH_CAP,
                max(60, int(_batch_img_cap * _per_img_timeout * 1.5)),
            )

            self.logger.info(
                "Processing %d pages, batch_size=%d, streaming=%s, text_pdf=%s",
                total_pages, batch_size,
                chunk_callback is not None,
                is_text_pdf,
            )

            return await self._process_pdf_batches(
                pdf_doc, pypdf_reader, total_pages, batch_size,
                image_count, team_config, auth_token,
                context_window, _batch_image_timeout,
                chunk_callback, all_chunks, is_text_pdf,
            )

        except Exception:  # noqa: BLE001
            self.logger.error(
                "Exception in extract_text_from_pdf_file", exc_info=True
            )
            return {"error": "Failed to extract text from PDF"}

        finally:
            if pdf_doc is not None:
                with suppress(Exception):
                    await loop.run_in_executor(_OCR_THREAD_EXECUTOR, pdf_doc.close)
            pypdf_reader = None
            if pdf_buffer is not None:
                with suppress(Exception):
                    pdf_buffer.close()
            gc.collect()

    async def _initialize_pdf_extraction(
        self,
        content: bytes,
        max_pages: Optional[int],
    ) -> Tuple[Any, Any, int]:
        """Open the PDF with pdfplumber, validate page count, return (doc, buffer, pages)."""
        if max_pages is None:
            max_pages = self.pdf_processing_config.get(
                "max_pages", ProcessingLimits.MAX_PDF_PAGES
            )

        pdf_buffer = _io.BytesIO(content)
        loop = asyncio.get_event_loop()

        self.logger.info("Opening PDF in thread executor to avoid blocking event loop")
        try:
            pdf_doc = await asyncio.wait_for(
                loop.run_in_executor(
                    _OCR_THREAD_EXECUTOR, lambda: pdfplumber.open(pdf_buffer)
                ),
                timeout=_PDF_OPEN_TIMEOUT,
            )
        except asyncio.TimeoutError:
            self.logger.error(
                "pdfplumber.open() timed out after %ds — file may be corrupt or adversarial",
                _PDF_OPEN_TIMEOUT,
            )
            return {"error": "PDF file could not be opened within the time limit. The file may be corrupt."}, None, 0  # type: ignore[return-value]

        total_pages = len(pdf_doc.pages)
        self.logger.info("PDF opened: %d total pages", total_pages)

        if max_pages and total_pages > max_pages:
            error_msg = (
                f"This PDF has {total_pages} pages, which exceeds the maximum "
                f"supported limit of {max_pages} pages per upload. "
                f"Please split the document into smaller files (each under "
                f"{max_pages} pages) and upload them separately."
            )
            self.logger.warning(
                "PDF rejected: %d pages exceeds cap of %d", total_pages, max_pages
            )
            return {"error": error_msg}, None, 0  # type: ignore[return-value]

        return pdf_doc, pdf_buffer, total_pages

    async def _process_pdf_batches(
        self,
        pdf_doc: Any,
        pypdf_reader: Optional[PdfReader],
        total_pages: int,
        batch_size: int,
        image_count: int,
        team_config: Dict,
        auth_token: str,
        context_window: int,
        batch_image_timeout: int,
        chunk_callback: Any,
        all_chunks: List,
        is_text_pdf: bool,
    ) -> Dict[str, Any]:
        """Iterate over page batches, extract content, stream or accumulate chunks."""
        for batch_start in range(0, total_pages, batch_size):
            gc.collect()
            batch_end = min(batch_start + batch_size, total_pages)

            batch_chunks, image_count = await self._process_pdf_page_batch(
                pdf_doc, pypdf_reader, batch_start, batch_end,
                total_pages, image_count, team_config, auth_token,
                context_window, batch_image_timeout, is_text_pdf,
            )

            if not batch_chunks:
                continue

            if chunk_callback is not None:
                await chunk_callback(batch_chunks)
            else:
                all_chunks.extend(batch_chunks)

        if chunk_callback is not None:
            success_msg = self.success_messages.get(
                "pdf_extracted", "Successfully extracted text from PDF file"
            )
            self.logger.info("%s (streaming mode)", success_msg)
            return {"chunks": [], "streamed": True}

        if not all_chunks:
            error_msg = self.error_messages.get(
                "pdf_no_content", "PDF contains no extractable content"
            )
            self.logger.error("%s", error_msg)
            return {"status_code": 400, "message": error_msg, "data": None}

        success_msg = self.success_messages.get(
            "pdf_extracted", "Successfully extracted text from PDF file"
        )
        self.logger.info("%s — %d total chunks", success_msg, len(all_chunks))
        return {
            "chunks": [
                {"text": c.page_content, **c.metadata} for c in all_chunks
            ]
        }

    async def _process_pdf_page_batch(
        self,
        pdf_doc: Any,
        pypdf_reader: Optional[PdfReader],
        batch_start: int,
        batch_end: int,
        total_pages: int,
        image_count: int,
        team_config: Dict,
        auth_token: str,
        context_window: int,
        batch_image_timeout: int,
        is_text_pdf: bool,
    ) -> Tuple[List, int]:
        """
        Process a batch of PDF pages: text (when text_pdf=True) + images.

        *is_text_pdf=False* skips the pdfplumber text path entirely, avoiding
        the 60-second stall on scanned / image-only PDFs.
        """
        batch_elements: List[Dict] = []
        batch_image_tasks: List = []
        batch_image_metadata: List = []
        loop = asyncio.get_event_loop()

        self.logger.info(
            "Processing pages %d-%d / %d",
            batch_start + 1, batch_end, total_pages,
        )

        for page_num in range(batch_start + 1, batch_end + 1):
            page = pdf_doc.pages[page_num - 1]
            page_text_content = ""
            try:
                if is_text_pdf:
                    try:
                        page_text_content = await asyncio.wait_for(
                            loop.run_in_executor(
                                _OCR_THREAD_EXECUTOR,
                                lambda p=page, pn=page_num: self._process_pdf_page_text(
                                    p, pn, batch_elements
                                ),
                            ),
                            timeout=60,
                        )
                    except asyncio.TimeoutError:
                        self.logger.warning(
                            "Page %d text extraction timed out after 60s — skipping",
                            page_num,
                        )

                if (
                    pypdf_reader is not None
                    and image_count < ProcessingLimits.MAX_PDF_IMAGES
                ):
                    page_tasks, page_meta, image_count = self._collect_pdf_page_images(
                        pypdf_reader, page_num, page_text_content,
                        context_window, image_count, team_config, auth_token,
                    )
                    batch_image_tasks.extend(page_tasks)
                    batch_image_metadata.extend(page_meta)
                elif image_count >= ProcessingLimits.MAX_PDF_IMAGES:
                    self.logger.warning(
                        "Reached max image limit (%d) for PDF",
                        ProcessingLimits.MAX_PDF_IMAGES,
                    )

            finally:
                with suppress(Exception):
                    del page

        if batch_image_tasks:
            await self._handle_batch_image_analysis(
                batch_image_tasks, batch_image_metadata, batch_elements,
                batch_start, batch_end, batch_image_timeout,
            )

        batch_image_tasks.clear()
        batch_image_metadata.clear()
        gc.collect()

        if not batch_elements:
            return [], image_count

        batch_chunks = await self._chunk_elements(batch_elements, source="pdf_file")
        batch_elements.clear()
        gc.collect()
        self.logger.info(
            "Batch pages %d-%d → %d chunks",
            batch_start + 1, batch_end, len(batch_chunks),
        )
        return batch_chunks, image_count

    async def _handle_batch_image_analysis(
        self,
        tasks: List,
        metadata: List,
        elements: List,
        batch_start: int,
        batch_end: int,
        timeout_sec: int,
    ) -> None:
        """
        Run image analysis for a batch with a hard timeout.

        On TimeoutError all pending futures are cancelled and awaited so their
        LLM-response buffers are freed immediately rather than leaking to the
        background.
        """
        futures = [asyncio.ensure_future(t) for t in tasks]
        try:
            async with asyncio.timeout(timeout_sec):
                await self._process_image_results(futures, metadata, elements, "pdf")
        except asyncio.TimeoutError:
            cancelled = sum(1 for f in futures if not f.done() and f.cancel())
            await asyncio.gather(*futures, return_exceptions=True)
            self.logger.error(
                "PDF image batch (pages %d-%d) timed out after %ds "
                "— cancelled %d pending task(s), skipped %d total image(s)",
                batch_start + 1, batch_end, timeout_sec,
                cancelled, len(tasks),
            )
        except Exception:  # noqa: BLE001
            for fut in futures:
                if not fut.done():
                    fut.cancel()
            await asyncio.gather(*futures, return_exceptions=True)
            self.logger.error(
                "Unexpected error in image batch (pages %d-%d)",
                batch_start + 1, batch_end,
                exc_info=True,
            )
            raise
        finally:
            futures.clear()
            gc.collect()

    def _process_pdf_page_text(
        self,
        page: Any,
        page_num: int,
        elements: List[Dict],
    ) -> str:
        """
        Extract text and tables from a single PDF page.

        Called via run_in_executor — must not access the event loop.
        """
        text_buffer: List[str] = []
        page_text_content: str = page.extract_text() or ""

        self.logger.debug(
            "Page %d: extracted %d chars", page_num, len(page_text_content)
        )

        if page_text_content:
            text_buffer.append(page_text_content)

        try:
            tables = page.extract_tables()
            for table_idx, table_data in enumerate(tables, 1):
                if text_buffer:
                    elements.append({
                        "type": "text",
                        "content": "\n".join(text_buffer),
                        "metadata": {"page_number": page_num},
                    })
                    text_buffer = []
                if table_data:
                    normalized = [
                        [cell if cell is not None else "" for cell in row]
                        for row in table_data
                    ]
                    elements.append({
                        "type": "table",
                        "content": normalized,
                        "metadata": {
                            "page_number": page_num,
                            "table_number": table_idx,
                        },
                    })
        except Exception:  # noqa: BLE001
            self.logger.warning(
                "Failed to extract tables from page %d", page_num, exc_info=True
            )

        if text_buffer:
            elements.append({
                "type": "text",
                "content": "\n".join(text_buffer),
                "metadata": {"page_number": page_num},
            })

        return page_text_content

    def _collect_pdf_page_images(
        self,
        pypdf_reader: PdfReader,
        page_num: int,
        page_text_content: str,
        context_window: int,
        image_count: int,
        team_config: Dict,
        auth_token: str,
    ) -> Tuple[List, List, int]:
        """Collect image-analysis tasks for one PDF page using pypdf xref."""
        image_tasks: List = []
        image_metadata: List = []
        page_image_count = 0

        try:
            pypdf_page = pypdf_reader.pages[page_num - 1]
            images_on_page = list(pypdf_page.images)
        except Exception:  # noqa: BLE001
            self.logger.warning(
                "pypdf could not access images on page %d", page_num, exc_info=True
            )
            return image_tasks, image_metadata, image_count

        for img_obj in images_on_page:
            if image_count >= ProcessingLimits.MAX_PDF_IMAGES:
                break
            if page_image_count >= ProcessingLimits.MAX_PDF_IMAGES_PER_PAGE:
                self.logger.debug(
                    "Page %d: per-page image cap (%d) reached",
                    page_num,
                    ProcessingLimits.MAX_PDF_IMAGES_PER_PAGE,
                )
                break

            img_index = page_image_count + 1
            task, meta = self._extract_pdf_image(
                img_obj, page_num, img_index,
                page_text_content, context_window, team_config, auth_token,
            )
            if task is not None:
                image_tasks.append(task)
                image_metadata.append(meta)
                image_count += 1
                page_image_count += 1

        return image_tasks, image_metadata, image_count

    def _extract_pdf_image(
        self,
        img_obj: Any,
        page_num: int,
        img_index: int,
        page_text_content: str,
        context_window: int,
        team_config: Dict,
        auth_token: str,
    ) -> Tuple[Optional[Any], Optional[Dict]]:
        """Extract raw bytes from a pypdf ImageFile and queue an analysis task."""
        try:
            image_bytes: bytes = img_obj.data
            if not image_bytes:
                return None, None

            width: int = getattr(img_obj, "width", 0) or 0
            height: int = getattr(img_obj, "height", 0) or 0

            if not self._is_relevant_image(
                image_bytes,
                width or None,
                height or None,
            ):
                self.logger.debug(
                    "Skipping small/decorative image on page %d, image %d",
                    page_num, img_index,
                )
                return None, None

            meta = {
                "page_number": page_num,
                "image_number": img_index,
                "image_name": getattr(img_obj, "name", f"img_{page_num}_{img_index}"),
            }

            task = self._analyze_image_with_semaphore(
                content=image_bytes,
                team_config=team_config,
                auth_token=auth_token,
                context=(
                    page_text_content[-context_window:]
                    if page_text_content else ""
                ),
            )
            image_bytes = b""  # release
            return task, meta

        except Exception:  # noqa: BLE001
            self.logger.warning(
                "Failed to extract image %d on page %d",
                img_index, page_num, exc_info=True,
            )
            return None, None

    # ── CSV extraction ────────────────────────────────────────────────────────

    def extract_text_from_csv(
        self, csv_bytes: bytes, filename: str = "file.csv"
    ) -> Dict[str, Any]:
        """Extract text content from CSV files with table preservation."""
        import csv
        from io import StringIO

        elements: List[Dict] = []
        try:
            decoded_text = self._decode_csv_bytes(csv_bytes)
            csv_bytes = b""  # release raw bytes
            csv_reader = csv.reader(StringIO(decoded_text))
            table_data = [row for row in csv_reader if row]

            if not table_data or len(table_data) <= 1:
                error_msg = self.error_messages.get(
                    "csv_no_content", "CSV file contains no data"
                )
                self.logger.error("%s", error_msg)
                return {"status_code": 400, "message": error_msg, "data": None}

            elements = [{"type": "table", "content": table_data, "metadata": {}}]
            chunks = self._run_async_chunking(
                elements, source=filename, use_semantic=False
            )

            success_msg = self.success_messages.get(
                "csv_extracted",
                "Successfully extracted {row_count} rows from CSV file",
            )
            self.logger.info("%s", success_msg.format(row_count=len(chunks)))
            return {"chunks": [{"text": d.page_content, **d.metadata} for d in chunks]}

        except Exception:  # noqa: BLE001
            self.logger.error("Exception during CSV text extraction", exc_info=True)
            return {"error": "Failed to extract text from CSV file"}

        finally:
            elements.clear()

    def _decode_csv_bytes(self, csv_bytes: bytes) -> str:
        """Try multiple encodings to decode CSV bytes."""
        encodings = self.text_extraction_config.get(
            "encodings",
            ["utf-8", "utf-16", "ascii", "latin-1", "iso-8859-1", "cp1252"],
        )
        for encoding in encodings:
            try:
                text = csv_bytes.decode(encoding)
                self.logger.info("CSV decoded using %s", encoding)
                return text
            except (UnicodeDecodeError, LookupError):
                continue

        fallback = self.text_extraction_config.get("fallback_encoding", "utf-8")
        text = csv_bytes.decode(fallback, errors="replace")
        self.logger.info("CSV decoded using %s (with replacements)", fallback)
        return text

    # ── Plain-text extraction ─────────────────────────────────────────────────

    def extract_text_from_text_file(self, content: bytes) -> str:
        """Decode plain text file bytes using a cascade of encodings."""
        if not isinstance(content, bytes):
            raise TypeError(f"Content must be bytes, got {type(content)}")
        if not content:
            self.logger.warning("Empty content provided to extract_text_from_text_file")
            return ""

        encodings = self.text_extraction_config.get(
            "encodings",
            ["utf-8", "utf-16", "ascii", "latin-1", "iso-8859-1", "cp1252"],
        )
        for encoding in encodings:
            try:
                text = content.decode(encoding, errors="strict")
                self.logger.info("Text file decoded using %s", encoding)
                return text
            except (UnicodeDecodeError, LookupError) as exc:
                self.logger.debug("Failed to decode with %s: %s", encoding, exc)

        fallback = self.text_extraction_config.get("fallback_encoding", "utf-8")
        errors = self.text_extraction_config.get("fallback_encoding_errors", "replace")
        try:
            text = content.decode(fallback, errors=errors)
            self.logger.info("Decoded with %s (errors=%s)", fallback, errors)
            return text
        except Exception as exc:
            self.logger.error("Failed to decode text file: %s", exc)
            raise ValueError(f"Could not decode text file: {exc}") from exc

    # ── Image extraction (standalone) ─────────────────────────────────────────

    async def extract_text_from_image(
        self,
        content: bytes,
        team_config: Dict[str, Any],
        auth_token: str,
    ) -> Dict[str, str]:
        """Extract text and understanding from a standalone image."""
        if not content:
            no_content_msg = self.error_messages.get(
                "no_content_extracted", "No content extracted from image"
            )
            self.logger.warning(
                "extract_text_from_image called with empty/None content"
            )
            return {"text": no_content_msg}

        try:
            analysis_result = await self._analyze_image_with_semaphore(
                content=content,
                team_config=team_config,
                auth_token=auth_token,
                context="",
            )
            formatted_text = self._format_image_analysis(analysis_result, "Image")
            no_content_msg = self.error_messages.get(
                "no_content_extracted", "No content extracted from image"
            )
            return {"text": formatted_text if formatted_text else no_content_msg}
        except Exception:  # noqa: BLE001
            self.logger.error(
                "Failed to extract text from image", exc_info=True
            )
            raise

    # ── Core image analysis ───────────────────────────────────────────────────

    async def _analyze_image_with_semaphore(
        self,
        content: bytes,
        team_config: Dict[str, Any],
        auth_token: str,
        context: str = "",
    ) -> Dict[str, Any]:
        """Semaphore-guarded wrapper around _analyze_image_for_understanding."""
        async with self._image_semaphore:
            return await self._analyze_image_for_understanding(
                content=content,
                team_config=team_config,
                auth_token=auth_token,
                context=context,
            )

    async def _analyze_image_for_understanding(
        self,
        content: bytes,
        team_config: Dict[str, Any],
        auth_token: str,
        context: str = "",
    ) -> Dict[str, Any]:
        """
        Comprehensive image analysis: OCR + semantic understanding via LLM.

        Memory notes:
        1. Raw bytes → base64 string; raw bytes released immediately.
        2. base64 string embedded into messages payload then released.
        3. messages cleared in finally.
        """
        image_base64: Optional[str] = None
        combined_messages: Optional[List] = None

        try:
            if not content:
                raise ValueError("Image content is empty or None")

            image_base64 = base64.b64encode(content).decode("utf-8")
            content = b""  # release raw bytes

            understanding_prompt = self._build_image_analysis_prompt(context)
            combined_messages = self._build_image_messages(
                image_base64, understanding_prompt
            )
            image_base64 = None  # embedded in messages; release standalone ref

            # Cap image analysis tokens to prevent unbounded 60K+ token responses
            # that caused 3-4 min LLM calls per image.  max_tokens is passed
            # via llm_params so it reaches litellm.acompletion correctly.
            # We merge into a shallow copy so the caller's team_config is unchanged.
            image_llm_params = {**team_config, "max_tokens": 2000}

            combined_text = await litellm_client.generate_response(
                llm_params=image_llm_params,
                messages=combined_messages,
                auth_token=auth_token,
            )
            self.logger.info(
                "Image analysis LLM call completed — response length: %d chars",
                len(combined_text),
            )

            result = self._parse_image_analysis_response(combined_text)
            success_msg = self.success_messages.get(
                "image_analysis_completed", "Comprehensive image analysis completed"
            )
            self.logger.info("%s", success_msg)
            return result

        except Exception:  # noqa: BLE001
            self.logger.error(
                "Failed to analyze image for understanding", exc_info=True
            )
            raise
        finally:
            image_base64 = None
            combined_messages = None

    def _build_image_analysis_prompt(self, context: str) -> str:
        """Construct the combined OCR + understanding prompt."""
        understanding_prompt_template: str = self.image_analysis_config.get(
            "understanding_prompt",
            "Analyze this image comprehensively and provide visual description, "
            "key insights, and data points.",
        )
        max_context_length = self.image_analysis_config.get("max_context_length", 200)
        context_included = (
            f"\nDocument Context: {context[:max_context_length]}" if context else ""
        )
        ocr_prefix = (
            "First, extract all text visible in this image under the section "
            "'OCR Text'.\nThen, analyze this image comprehensively following the "
            "instructions below.\n\n"
        )
        prompt = f"{ocr_prefix}{understanding_prompt_template}{context_included}"
        return prompt.format(reasoning_section=REASONING_SECTION_PROMPT)

    @staticmethod
    def _build_image_messages(
        image_base64: str, prompt: str
    ) -> List[Dict[str, Any]]:
        """Build the messages list for the image-analysis API call."""
        return [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/png;base64,{image_base64}"
                        },
                    },
                ],
            }
        ]

    def _parse_image_analysis_response(
        self, combined_text: str
    ) -> Dict[str, Any]:
        """Parse the combined image analysis LLM response into structured sections."""
        result: Dict[str, Any] = {
            "ocr_text": None,
            "visual_description": None,
            "key_insights": None,
            "data_points": None,
        }

        keyword_map = self._get_image_section_keyword_map()
        current_section: Optional[str] = None
        section_content: List[str] = []

        for line in combined_text.split("\n"):
            matched = self._match_image_section(line.lower(), keyword_map)
            if matched:
                if current_section and section_content:
                    result[current_section] = "\n".join(section_content)
                current_section = matched
                section_content = []
            elif current_section and line:
                section_content.append(line)

        if current_section and section_content:
            result[current_section] = "\n".join(section_content)

        if not any(
            result[k]
            for k in ("visual_description", "key_insights", "data_points")
        ):
            result["visual_description"] = combined_text

        return result

    def _get_image_section_keyword_map(self) -> Dict[str, List[str]]:
        """Return the keyword map used to detect section headers in LLM output."""
        return {
            "ocr_text": self.image_analysis_config.get(
                "ocr_keywords", ["ocr_text", "ocr text", "extracted text"]
            ),
            "visual_description": self.image_analysis_config.get(
                "visual_description_keywords",
                ["visual_description", "visual description", "image type"],
            ),
            "key_insights": self.image_analysis_config.get(
                "key_insights_keywords",
                ["key_insights", "key insights", "key_insight", "main messages", "findings"],
            ),
            "data_points": self.image_analysis_config.get(
                "data_points_keywords",
                ["data_points", "data points", "data_point", "metrics", "numbers"],
            ),
        }

    @staticmethod
    def _match_image_section(
        line_lower: str,
        keyword_map: Dict[str, List[str]],
    ) -> Optional[str]:
        """Return the matched section name if *line_lower* contains any keyword."""
        for section, keywords in keyword_map.items():
            if any(kw in line_lower for kw in keywords):
                return section
        return None

    # ── ZIP extraction ────────────────────────────────────────────────────────

    async def extract_text_from_zip(
        self,
        content: bytes,
        team_config: Dict[str, Any],
        auth_token: str,
        current_depth: int = 0,
    ) -> Dict[str, Any]:
        """Extract text from all supported files within a zip archive."""
        if current_depth >= ProcessingLimits.MAX_ZIP_DEPTH:
            error_msg = (
                f"Maximum ZIP nesting depth "
                f"({ProcessingLimits.MAX_ZIP_DEPTH}) exceeded"
            )
            self.logger.error("%s", error_msg)
            return {"error": error_msg}

        extracted_files: Dict[str, Any] = {}
        zip_buffer: Optional[BytesIO] = None

        try:
            zip_buffer = BytesIO(content)
            content = b""  # release raw bytes
            supported_text_ext, supported_image_ext = (
                self._get_zip_supported_extensions()
            )

            with zipfile.ZipFile(zip_buffer, "r") as zip_file:
                total_size = 0
                files_processed = 0
                for file_name in zip_file.namelist():
                    if files_processed >= ProcessingLimits.MAX_ZIP_FILES:
                        self.logger.warning(
                            "Reached maximum file limit (%d) in ZIP",
                            ProcessingLimits.MAX_ZIP_FILES,
                        )
                        break
                    if file_name.endswith("/"):
                        continue

                    should_stop, total_size, files_processed = (
                        await self._process_zip_entry(
                            zip_file, file_name, extracted_files,
                            total_size, files_processed,
                            supported_text_ext, supported_image_ext,
                            team_config, auth_token, current_depth,
                        )
                    )
                    if should_stop:
                        break

            success_msg = self.success_messages.get(
                "zip_extracted",
                "Successfully extracted text from {count} files in zip archive",
            )
            self.logger.info("%s", success_msg.format(count=len(extracted_files)))
            return extracted_files

        except zipfile.BadZipFile as exc:
            error_msg = self.error_messages.get(
                "zip_invalid", "Invalid zip file: {error}"
            )
            self.logger.error("%s", error_msg.format(error=str(exc)))
            return {"error": error_msg.format(error=str(exc))}

        except Exception:  # noqa: BLE001
            self.logger.error("Exception during zip extraction", exc_info=True)
            return {"error": "Failed to extract zip file"}

        finally:
            if zip_buffer is not None:
                with suppress(Exception):
                    zip_buffer.close()

    def _get_zip_supported_extensions(self) -> Tuple[List[str], List[str]]:
        """Return (text_extensions, image_extensions) for zip dispatch."""
        text_ext = self.text_extraction_config.get(
            "supported_extensions",
            ["txt", "md", "csv", "json", "xml", "html", "py", "js",
             "java", "cpp", "c", "h"],
        )
        image_ext = self.text_extraction_config.get(
            "image_extensions",
            ["png", "jpg", "jpeg", "gif", "bmp", "tiff", "webp"],
        )
        return text_ext, image_ext

    async def _process_zip_entry(
        self,
        zip_file: zipfile.ZipFile,
        file_name: str,
        extracted_files: Dict,
        total_extracted_size: int,
        files_processed: int,
        supported_text_ext: List[str],
        supported_image_ext: List[str],
        team_config: Dict,
        auth_token: str,
        current_depth: int,
    ) -> Tuple[bool, int, int]:
        """
        Process one entry in the ZIP; return (should_stop, total_size, count).
        gc.collect() is called after each entry to reclaim large document memory.
        """
        file_bytes: Optional[bytes] = None
        try:
            file_bytes = zip_file.read(file_name)
            file_size = len(file_bytes)

            if total_extracted_size + file_size > ProcessingLimits.MAX_ZIP_EXTRACTED_SIZE:
                self.logger.warning(
                    "Reached maximum extracted size limit (%d bytes) in ZIP",
                    ProcessingLimits.MAX_ZIP_EXTRACTED_SIZE,
                )
                return True, total_extracted_size, files_processed

            total_extracted_size += file_size
            files_processed += 1
            file_ext = file_name.lower().rsplit(".", 1)[-1] if "." in file_name else ""

            await self._dispatch_zip_file(
                file_ext, file_name, file_bytes, extracted_files,
                supported_text_ext, supported_image_ext,
                team_config, auth_token, current_depth,
            )

        except Exception:  # noqa: BLE001
            self.logger.error(
                "Failed to process %s in ZIP", file_name, exc_info=True
            )
            extracted_files[file_name] = [
                {"text": "Error processing file", "type": "error"}
            ]
        finally:
            file_bytes = None
            gc.collect()

        return False, total_extracted_size, files_processed

    async def _dispatch_zip_file(
        self,
        file_ext: str,
        file_name: str,
        file_bytes: bytes,
        extracted_files: Dict,
        supported_text_ext: List[str],
        supported_image_ext: List[str],
        team_config: Dict,
        auth_token: str,
        current_depth: int,
    ) -> None:
        """Route a ZIP entry to the correct extractor based on file extension."""
        if file_ext == "zip":
            await self._handle_nested_zip(
                file_name, file_bytes, extracted_files,
                team_config, auth_token, current_depth,
            )
        elif file_ext in ("docx", "doc"):
            result = await self.extract_text_from_docx(
                file_bytes, team_config, auth_token
            )
            self._store_zip_result(file_name, extracted_files, result)
        elif file_ext in ("xlsx", "xls"):
            result = await self.extract_text_from_excel(file_bytes, file_name)
            self._store_zip_result(file_name, extracted_files, result)
        elif file_ext in ("ppt", "pptx"):
            result = await self.extract_text_from_ppt(
                file_bytes, team_config, auth_token
            )
            self._store_zip_result(file_name, extracted_files, result)
        elif file_ext == "pdf":
            result = await self.extract_text_from_pdf_file(
                file_bytes, team_config, auth_token
            )
            self._store_zip_result(file_name, extracted_files, result)
        elif file_ext == "csv":
            result = self.extract_text_from_csv(file_bytes, file_name)
            self._store_zip_result(file_name, extracted_files, result)
        elif file_ext in supported_text_ext:
            self._handle_zip_text_file(file_name, file_bytes, extracted_files)
        elif file_ext in supported_image_ext:
            await self._handle_zip_image_file(
                file_name, file_bytes, extracted_files, team_config, auth_token
            )
        else:
            unsupported_msg = self.error_messages.get(
                "unsupported_file_type", "Unsupported file type: .{file_ext}"
            )
            self.logger.warning("%s", unsupported_msg.format(file_ext=file_ext))
            extracted_files[file_name] = [
                {"text": unsupported_msg.format(file_ext=file_ext), "type": "unsupported"}
            ]

    async def _handle_nested_zip(
        self,
        file_name: str,
        file_bytes: bytes,
        extracted_files: Dict,
        team_config: Dict,
        auth_token: str,
        current_depth: int,
    ) -> None:
        """Handle a nested ZIP file inside a ZIP."""
        self.logger.info("Processing nested zip: %s", file_name)
        nested_results = await self.extract_text_from_zip(
            file_bytes, team_config, auth_token,
            current_depth=current_depth + 1,
        )
        if isinstance(nested_results, dict) and "error" not in nested_results:
            for nested_name, nested_chunks in nested_results.items():
                extracted_files[f"{file_name}/{nested_name}"] = nested_chunks
        else:
            error_text = (
                nested_results.get("error", "Unknown error")
                if isinstance(nested_results, dict)
                else "Unknown error"
            )
            extracted_files[file_name] = [
                {"text": f"Error: {error_text}", "type": "error"}
            ]

    @staticmethod
    def _store_zip_result(
        file_name: str,
        extracted_files: Dict,
        result: Dict[str, Any],
    ) -> None:
        """
        Store the extraction result for a ZIP entry.

        Renamed from _handle_zip_result (was misleadingly named and called
        with `await` in some places despite being sync).
        """
        if "chunks" in result:
            extracted_files[file_name] = result["chunks"]
        else:
            error_text = result.get("error", "Unknown error")
            extracted_files[file_name] = [
                {"text": f"Error: {error_text}", "type": "error"}
            ]

    def _handle_zip_text_file(
        self,
        file_name: str,
        file_bytes: bytes,
        extracted_files: Dict,
    ) -> None:
        """Decode and store a plain-text file from inside a ZIP."""
        try:
            text = self.extract_text_from_text_file(file_bytes)
            extracted_files[file_name] = [{"text": text, "type": "text"}]
        except Exception:  # noqa: BLE001
            self.logger.error(
                "Failed to decode text file %s in ZIP", file_name, exc_info=True
            )
            extracted_files[file_name] = [
                {"text": "Error: could not decode file", "type": "error"}
            ]

    async def _handle_zip_image_file(
        self,
        file_name: str,
        file_bytes: bytes,
        extracted_files: Dict,
        team_config: Dict,
        auth_token: str,
    ) -> None:
        """Analyse and store an image file from inside a ZIP."""
        try:
            result = await self._analyze_image_with_semaphore(
                file_bytes, team_config, auth_token=auth_token, context=""
            )
            image_output = self._format_image_analysis(result, file_name)
            extracted_files[file_name] = [
                {
                    "text": image_output if image_output else "No content extracted",
                    "type": "image",
                }
            ]
        except Exception:  # noqa: BLE001
            self.logger.error(
                "Failed to analyse image %s in ZIP", file_name, exc_info=True
            )
            extracted_files[file_name] = [
                {"text": "Error: image analysis failed", "type": "error"}
            ]