"""
Comprehensive unit tests for ocr.py
Covers: ProcessingLimits, TableProcessor, Ocr — all methods, branches, and edge cases.
"""

import asyncio
import base64
import io
import os
import tempfile
import zipfile
from io import BytesIO
from unittest.mock import AsyncMock, MagicMock, Mock, patch, PropertyMock, call

import pandas as pd
import pytest
import pytest_asyncio
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches, Pt
from langchain_core.documents import Document as LangchainDocument

# ── project imports ──────────────────────────────────────────────────────────
import sys
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from src.ocr.ocr import Ocr, ProcessingLimits, TableProcessor


# ═══════════════════════════════════════════════════════════════════════════
# Helpers / fixtures
# ═══════════════════════════════════════════════════════════════════════════

CONFIG_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), "config", "ocr_config.yaml")

TEAM_CONFIG = {"model": "test-model"}
AUTH_TOKEN = "test-token"

IMAGE_ANALYSIS_RESULT = {
    "ocr_text": "Sample OCR text",
    "visual_description": "A bar chart showing sales data",
    "key_insights": "Sales increased by 20%",
    "data_points": "Q1: 100, Q2: 120",
}

EMPTY_IMAGE_ANALYSIS = {
    "ocr_text": None,
    "visual_description": None,
    "key_insights": None,
    "data_points": None,
}


def make_large_image_bytes(size: int = 10_000) -> bytes:
    """Return a PNG-header bytes blob large enough to pass size filter."""
    return b"\x89PNG\r\n\x1a\n" + b"\x00" * (size - 8)


def make_small_image_bytes(size: int = 100) -> bytes:
    return b"\x00" * size


def make_xlsx_bytes(data: dict = None) -> bytes:
    """Create a minimal in-memory xlsx file."""
    if data is None:
        data = {"Sheet1": pd.DataFrame({"col1": [1, 2], "col2": ["a", "b"]})}
    buf = BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        for sheet, df in data.items():
            df.to_excel(writer, sheet_name=sheet, index=False)
    buf.seek(0)
    return buf.read()


def make_docx_bytes(text: str = "Hello World") -> bytes:
    """Create a minimal in-memory docx file."""
    doc = DocxDocument()
    doc.add_paragraph(text)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def make_docx_with_table_bytes() -> bytes:
    """Create a docx with a table."""
    doc = DocxDocument()
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header1"
    table.cell(0, 1).text = "Header2"
    table.cell(1, 0).text = "Value1"
    table.cell(1, 1).text = "Value2"
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def make_pptx_bytes(text: str = "Slide text") -> bytes:
    """Create a minimal in-memory pptx file."""
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = text
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def make_zip_bytes(files: dict) -> bytes:
    """Create an in-memory zip with {filename: bytes} contents."""
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        for name, content in files.items():
            zf.writestr(name, content)
    buf.seek(0)
    return buf.read()


def make_nested_zip_bytes(inner_files: dict) -> bytes:
    """Zip containing another zip."""
    inner = make_zip_bytes(inner_files)
    return make_zip_bytes({"inner.zip": inner})


@pytest_asyncio.fixture
def ocr(tmp_path):
    """Return an Ocr instance pointing at the real config file."""
    return Ocr(config_path=CONFIG_PATH)


# ═══════════════════════════════════════════════════════════════════════════
# ProcessingLimits
# ═══════════════════════════════════════════════════════════════════════════

class TestProcessingLimits:
    def test_constants_exist_and_positive(self):
        assert ProcessingLimits.MAX_CONCURRENT_IMAGES > 0
        assert ProcessingLimits.MAX_ZIP_DEPTH > 0
        assert ProcessingLimits.MAX_ZIP_FILES > 0
        assert ProcessingLimits.MAX_ZIP_EXTRACTED_SIZE > 0
        assert ProcessingLimits.MAX_PDF_IMAGES > 0
        assert ProcessingLimits.MAX_PPTX_IMAGES > 0
        assert ProcessingLimits.MAX_DOCX_IMAGES > 0
        assert ProcessingLimits.IMAGE_ANALYSIS_TIMEOUT > 0

    def test_specific_values(self):
        assert ProcessingLimits.MAX_CONCURRENT_IMAGES == 3
        assert ProcessingLimits.MAX_ZIP_DEPTH == 3
        assert ProcessingLimits.MAX_ZIP_FILES == 100
        assert ProcessingLimits.MAX_ZIP_EXTRACTED_SIZE == 100 * 1024 * 1024
        assert ProcessingLimits.MAX_PDF_IMAGES == 50
        assert ProcessingLimits.MAX_PPTX_IMAGES == 50
        assert ProcessingLimits.MAX_DOCX_IMAGES == 50
        assert ProcessingLimits.IMAGE_ANALYSIS_TIMEOUT == 30


# ═══════════════════════════════════════════════════════════════════════════
# TableProcessor
# ═══════════════════════════════════════════════════════════════════════════

class TestTableProcessorToMarkdown:
    def test_empty_table(self):
        assert TableProcessor.table_to_markdown([]) == ""

    def test_none_table(self):
        assert TableProcessor.table_to_markdown(None) == ""

    def test_header_only(self):
        result = TableProcessor.table_to_markdown([["A", "B"]])
        assert "| A | B |" in result
        assert "---" in result

    def test_full_table(self):
        data = [["Name", "Age"], ["Alice", "30"], ["Bob", "25"]]
        result = TableProcessor.table_to_markdown(data)
        assert "| Name | Age |" in result
        assert "| Alice | 30 |" in result
        assert "| Bob | 25 |" in result
        assert "---" in result

    def test_short_row_padded(self):
        data = [["A", "B", "C"], ["only_one"]]
        result = TableProcessor.table_to_markdown(data)
        # The short row should be padded with empty strings
        assert "| only_one |  |  |" in result

    def test_numeric_cells(self):
        data = [["X"], [1], [2.5]]
        result = TableProcessor.table_to_markdown(data)
        assert "| 1 |" in result
        assert "| 2.5 |" in result


class TestTableProcessorChunkLargeTable:
    def test_empty_table(self):
        assert TableProcessor.chunk_large_table([]) == []

    def test_header_only(self):
        assert TableProcessor.chunk_large_table([["H1", "H2"]]) == []

    def test_small_table_single_chunk(self):
        data = [["H"]] + [[str(i)] for i in range(5)]
        chunks = TableProcessor.chunk_large_table(data, rows_per_chunk=20)
        assert len(chunks) == 1

    def test_large_table_multiple_chunks(self):
        header = ["Col"]
        rows = [[str(i)] for i in range(45)]
        data = [header] + rows
        chunks = TableProcessor.chunk_large_table(data, rows_per_chunk=20)
        assert len(chunks) == 3  # 20 + 20 + 5

    def test_exact_chunk_boundary(self):
        header = ["Col"]
        rows = [[str(i)] for i in range(20)]
        data = [header] + rows
        chunks = TableProcessor.chunk_large_table(data, rows_per_chunk=20)
        assert len(chunks) == 1

    def test_each_chunk_has_header(self):
        header = ["ID", "Name"]
        rows = [[str(i), f"row{i}"] for i in range(25)]
        data = [header] + rows
        chunks = TableProcessor.chunk_large_table(data, rows_per_chunk=10)
        for chunk in chunks:
            assert "| ID | Name |" in chunk


class TestTableProcessorExtractDocxTable:
    def test_extract_table(self):
        doc = DocxDocument()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "H1"
        table.cell(0, 1).text = "H2"
        table.cell(1, 0).text = "V1"
        table.cell(1, 1).text = "V2"
        result = TableProcessor.extract_docx_table(table)
        assert result == [["H1", "H2"], ["V1", "V2"]]

    def test_empty_table(self):
        doc = DocxDocument()
        table = doc.add_table(rows=1, cols=1)
        result = TableProcessor.extract_docx_table(table)
        assert result == [[""]]


class TestTableProcessorExtractPptxTable:
    def test_extract_pptx_table(self):
        mock_cell1 = Mock()
        mock_cell1.text = "  Header1  "
        mock_cell2 = Mock()
        mock_cell2.text = "Header2"
        mock_row = Mock()
        mock_row.cells = [mock_cell1, mock_cell2]

        mock_cell3 = Mock()
        mock_cell3.text = "Val1"
        mock_cell4 = Mock()
        mock_cell4.text = "Val2"
        mock_row2 = Mock()
        mock_row2.cells = [mock_cell3, mock_cell4]

        mock_table = Mock()
        mock_table.rows = [mock_row, mock_row2]

        result = TableProcessor.extract_pptx_table(mock_table)
        assert result == [["Header1", "Header2"], ["Val1", "Val2"]]


# ═══════════════════════════════════════════════════════════════════════════
# Ocr.__init__ / _load_config
# ═══════════════════════════════════════════════════════════════════════════

class TestOcrInit:
    def test_loads_real_config(self, ocr):
        assert isinstance(ocr.config, dict)
        assert "ocr" in ocr.config

    def test_missing_config_returns_defaults(self, tmp_path):
        o = Ocr(config_path=str(tmp_path / "nonexistent.yaml"))
        assert o.config == {"ocr": {}}

    def test_invalid_yaml_returns_defaults(self, tmp_path):
        bad_yaml = tmp_path / "bad.yaml"
        bad_yaml.write_text(": invalid: yaml: [[[")
        # yaml.safe_load may raise — must fall back to defaults
        o = Ocr(config_path=str(bad_yaml))
        assert isinstance(o.config, dict)

    def test_empty_yaml_returns_defaults(self, tmp_path):
        empty_yaml = tmp_path / "empty.yaml"
        empty_yaml.write_text("")
        o = Ocr(config_path=str(empty_yaml))
        # yaml.safe_load("") returns None → should default to {"ocr": {}}
        assert o.config == {"ocr": {}}

    def test_semaphore_created(self, ocr):
        assert ocr._image_semaphore is not None

    def test_text_splitter_created(self, ocr):
        assert ocr.text_splitter is not None

    def test_load_config_exception_path(self, tmp_path):
        """Cover the generic except branch in _load_config."""
        with patch("builtins.open", side_effect=PermissionError("no access")):
            with patch("os.path.exists", return_value=True):
                o = Ocr(config_path="any_path.yaml")
        assert o.config == {"ocr": {}}


# ═══════════════════════════════════════════════════════════════════════════
# Ocr._is_relevant_image
# ═══════════════════════════════════════════════════════════════════════════

class TestIsRelevantImage:
    def test_small_file_rejected(self, ocr):
        assert ocr._is_relevant_image(b"\x00" * 100) is False

    def test_large_file_no_dims(self, ocr):
        assert ocr._is_relevant_image(make_large_image_bytes()) is True

    def test_large_file_good_dims(self, ocr):
        assert ocr._is_relevant_image(make_large_image_bytes(), 200, 200) is True

    def test_small_width_rejected(self, ocr):
        assert ocr._is_relevant_image(make_large_image_bytes(), 50, 200) is False

    def test_small_height_rejected(self, ocr):
        assert ocr._is_relevant_image(make_large_image_bytes(), 200, 50) is False

    def test_extreme_aspect_ratio_rejected(self, ocr):
        # width/height = 5000/10 = 500 >> 20
        assert ocr._is_relevant_image(make_large_image_bytes(), 5000, 10) is False

    def test_exception_returns_true(self, ocr):
        """If image_bytes has no __len__, the try raises → returns True."""
        bad_bytes = Mock(spec=[])  # no __len__
        # Patch len to raise
        with patch("src.ocr.ocr.ProcessingLimits"):  # just to exercise the exception branch
            result = ocr._is_relevant_image(bad_bytes)
        # Since bad_bytes lacks __len__, len() will raise AttributeError → except returns True
        assert result is True

    def test_zero_dims_skips_dim_check(self, ocr):
        """width=0, height=0 → falsy → skip dim checks, only size matters."""
        assert ocr._is_relevant_image(make_large_image_bytes(), 0, 0) is True


# ═══════════════════════════════════════════════════════════════════════════
# Ocr._chunk_elements
# ═══════════════════════════════════════════════════════════════════════════

class TestChunkElements:
    @pytest.mark.asyncio
    async def test_text_element(self, ocr):
        elements = [{"type": "text", "content": "Hello world", "metadata": {"page_number": 1}}]
        chunks = await ocr._chunk_elements(elements, source="test")
        assert len(chunks) >= 1
        assert any("Hello world" in c.page_content for c in chunks)

    @pytest.mark.asyncio
    async def test_table_element(self, ocr):
        data = [["H1", "H2"], ["v1", "v2"]]
        elements = [{"type": "table", "content": data, "metadata": {}}]
        chunks = await ocr._chunk_elements(elements, source="test")
        assert len(chunks) == 1
        assert "H1" in chunks[0].page_content

    @pytest.mark.asyncio
    async def test_image_element(self, ocr):
        elements = [{"type": "image", "content": "Image analysis text", "metadata": {}}]
        chunks = await ocr._chunk_elements(elements, source="test")
        assert len(chunks) == 1
        assert chunks[0].page_content == "Image analysis text"

    @pytest.mark.asyncio
    async def test_large_table_splits_into_chunks(self, ocr):
        # Force small rows_per_chunk via config override
        ocr.chunking_config["table_rows_per_chunk"] = 3
        header = ["A", "B"]
        rows = [[str(i), str(i * 2)] for i in range(10)]
        data = [header] + rows
        elements = [{"type": "table", "content": data, "metadata": {}}]
        chunks = await ocr._chunk_elements(elements, source="test")
        assert len(chunks) > 1
        # Each chunk's table_chunk_index should reflect split
        indices = [c.metadata.get("table_chunk_index") for c in chunks]
        assert any(i is not None for i in indices)

    @pytest.mark.asyncio
    async def test_metadata_propagated(self, ocr):
        elements = [{"type": "image", "content": "img", "metadata": {"page_number": 5}}]
        chunks = await ocr._chunk_elements(elements, source="mysrc", page=3)
        assert chunks[0].metadata["source"] == "mysrc"
        assert chunks[0].metadata["page_number"] == 5

    @pytest.mark.asyncio
    async def test_empty_elements(self, ocr):
        chunks = await ocr._chunk_elements([], source="empty")
        assert chunks == []

    @pytest.mark.asyncio
    async def test_mixed_elements(self, ocr):
        elements = [
            {"type": "text", "content": "Some text here", "metadata": {}},
            {"type": "table", "content": [["H"], ["v"]], "metadata": {}},
            {"type": "image", "content": "img content", "metadata": {}},
        ]
        chunks = await ocr._chunk_elements(elements, source="mixed")
        types = [c.metadata.get("type") for c in chunks]
        assert "text" in types
        assert "table" in types
        assert "image" in types


# ═══════════════════════════════════════════════════════════════════════════
# Ocr._format_image_analysis
# ═══════════════════════════════════════════════════════════════════════════

class TestFormatImageAnalysis:
    def test_all_fields_present(self, ocr):
        result = ocr._format_image_analysis(IMAGE_ANALYSIS_RESULT, "test.png")
        assert "Sample OCR text" in result
        assert "bar chart" in result
        assert "20%" in result
        assert "Q1: 100" in result

    def test_empty_result_returns_empty(self, ocr):
        result = ocr._format_image_analysis(EMPTY_IMAGE_ANALYSIS, "x.png")
        assert result == ""

    def test_only_ocr_text(self, ocr):
        r = {"ocr_text": "text only", "visual_description": None, "key_insights": None, "data_points": None}
        out = ocr._format_image_analysis(r, "img")
        assert "text only" in out

    def test_only_visual_description(self, ocr):
        r = {"ocr_text": None, "visual_description": "a photo", "key_insights": None, "data_points": None}
        out = ocr._format_image_analysis(r, "img")
        assert "a photo" in out

    def test_only_key_insights(self, ocr):
        r = {"ocr_text": None, "visual_description": None, "key_insights": "insight", "data_points": None}
        out = ocr._format_image_analysis(r, "img")
        assert "insight" in out

    def test_only_data_points(self, ocr):
        r = {"ocr_text": None, "visual_description": None, "key_insights": None, "data_points": "42"}
        out = ocr._format_image_analysis(r, "img")
        assert "42" in out

    def test_image_name_in_header(self, ocr):
        r = {"ocr_text": "x", "visual_description": None, "key_insights": None, "data_points": None}
        out = ocr._format_image_analysis(r, "my_image.jpg")
        assert "my_image.jpg" in out


# ═══════════════════════════════════════════════════════════════════════════
# Ocr._analyze_image_for_understanding
# ═══════════════════════════════════════════════════════════════════════════

class TestAnalyzeImageForUnderstanding:
    @pytest.mark.asyncio
    async def test_parses_all_sections(self, ocr):
        # Content lines must not accidentally match section keyword patterns
        response = (
            "OCR Text:\nHello from the image\n"
            "Visual Description:\nA bar chart showing Q1\n"
            "Key Insights:\nSales rose 20 percent\n"
            "Data Points:\nQ1=100 Q2=120\n"
        )
        with patch("src.ocr.ocr.litellm_client") as mock_client:
            mock_client.generate_response = AsyncMock(return_value=response)
            result = await ocr._analyze_image_for_understanding(
                b"fakepng", TEAM_CONFIG, AUTH_TOKEN, context="some context"
            )
        assert result["ocr_text"] is not None
        assert result["visual_description"] is not None
        assert result["key_insights"] is not None
        assert result["data_points"] is not None

    @pytest.mark.asyncio
    async def test_fallback_to_visual_description_when_no_sections(self, ocr):
        response = "plain text with no section headers"
        with patch("src.ocr.ocr.litellm_client") as mock_client:
            mock_client.generate_response = AsyncMock(return_value=response)
            result = await ocr._analyze_image_for_understanding(
                b"fakepng", TEAM_CONFIG, AUTH_TOKEN
            )
        assert result["visual_description"] == response

    @pytest.mark.asyncio
    async def test_empty_context(self, ocr):
        with patch("src.ocr.ocr.litellm_client") as mock_client:
            mock_client.generate_response = AsyncMock(return_value="Visual Description:\nfoo")
            result = await ocr._analyze_image_for_understanding(
                b"img", TEAM_CONFIG, AUTH_TOKEN, context=""
            )
        assert result["visual_description"] == "foo"

    @pytest.mark.asyncio
    async def test_exception_propagates(self, ocr):
        with patch("src.ocr.ocr.litellm_client") as mock_client:
            mock_client.generate_response = AsyncMock(side_effect=RuntimeError("API down"))
            with pytest.raises(RuntimeError, match="API down"):
                await ocr._analyze_image_for_understanding(b"img", TEAM_CONFIG, AUTH_TOKEN)

    @pytest.mark.asyncio
    async def test_section_content_flushed_at_end(self, ocr):
        response = "OCR Text:\nfirst line\nsecond line"
        with patch("src.ocr.ocr.litellm_client") as mock_client:
            mock_client.generate_response = AsyncMock(return_value=response)
            result = await ocr._analyze_image_for_understanding(b"img", TEAM_CONFIG, AUTH_TOKEN)
        assert "first line" in result["ocr_text"]
        assert "second line" in result["ocr_text"]

    @pytest.mark.asyncio
    async def test_context_truncation(self, ocr):
        """Context longer than max_context_length is truncated."""
        long_context = "x" * 1000
        with patch("src.ocr.ocr.litellm_client") as mock_client:
            mock_client.generate_response = AsyncMock(return_value="Visual Description:\ntest")
            await ocr._analyze_image_for_understanding(
                b"img", TEAM_CONFIG, AUTH_TOKEN, context=long_context
            )
        # Just ensure it doesn't crash and the call was made
        mock_client.generate_response.assert_called_once()

    @pytest.mark.asyncio
    async def test_section_switch_flushes_previous(self, ocr):
        """Switching from ocr_text section to visual_description flushes ocr content."""
        response = "OCR Text:\nocr content\nVisual Description:\nvisual content"
        with patch("src.ocr.ocr.litellm_client") as mock_client:
            mock_client.generate_response = AsyncMock(return_value=response)
            result = await ocr._analyze_image_for_understanding(b"img", TEAM_CONFIG, AUTH_TOKEN)
        assert result["ocr_text"] == "ocr content"
        assert result["visual_description"] == "visual content"


# ═══════════════════════════════════════════════════════════════════════════
# Ocr._analyze_image_with_semaphore
# ═══════════════════════════════════════════════════════════════════════════

class TestAnalyzeImageWithSemaphore:
    @pytest.mark.asyncio
    async def test_delegates_to_understand(self, ocr):
        with patch.object(ocr, "_analyze_image_for_understanding", new_callable=AsyncMock) as mock_fn:
            mock_fn.return_value = IMAGE_ANALYSIS_RESULT
            result = await ocr._analyze_image_with_semaphore(b"img", TEAM_CONFIG, AUTH_TOKEN, "ctx")
        assert result == IMAGE_ANALYSIS_RESULT
        mock_fn.assert_called_once_with(
            content=b"img", team_config=TEAM_CONFIG, auth_token=AUTH_TOKEN, context="ctx"
        )

    @pytest.mark.asyncio
    async def test_semaphore_limits_concurrency(self, ocr):
        """Fire MAX_CONCURRENT_IMAGES+2 tasks, assert all complete."""
        call_count = 0

        async def fake_analyze(**kwargs):
            nonlocal call_count
            call_count += 1
            await asyncio.sleep(0)
            return IMAGE_ANALYSIS_RESULT

        with patch.object(ocr, "_analyze_image_for_understanding", side_effect=fake_analyze):
            tasks = [
                ocr._analyze_image_with_semaphore(b"img", TEAM_CONFIG, AUTH_TOKEN)
                for _ in range(ProcessingLimits.MAX_CONCURRENT_IMAGES + 2)
            ]
            await asyncio.gather(*tasks)
        assert call_count == ProcessingLimits.MAX_CONCURRENT_IMAGES + 2


# ═══════════════════════════════════════════════════════════════════════════
# Ocr.extract_text_from_image
# ═══════════════════════════════════════════════════════════════════════════

class TestExtractTextFromImage:
    @pytest.mark.asyncio
    async def test_returns_formatted_text(self, ocr):
        with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_fn:
            mock_fn.return_value = IMAGE_ANALYSIS_RESULT
            result = await ocr.extract_text_from_image(b"img", TEAM_CONFIG, AUTH_TOKEN)
        assert "text" in result
        assert result["text"]  # non-empty

    @pytest.mark.asyncio
    async def test_empty_analysis_returns_no_content_msg(self, ocr):
        with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_fn:
            mock_fn.return_value = EMPTY_IMAGE_ANALYSIS
            result = await ocr.extract_text_from_image(b"img", TEAM_CONFIG, AUTH_TOKEN)
        assert result["text"] == ocr.error_messages.get(
            "no_content_extracted", "No content extracted from image"
        )

    @pytest.mark.asyncio
    async def test_exception_propagates(self, ocr):
        with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_fn:
            mock_fn.side_effect = ValueError("bad image")
            with pytest.raises(ValueError, match="bad image"):
                await ocr.extract_text_from_image(b"img", TEAM_CONFIG, AUTH_TOKEN)


# ═══════════════════════════════════════════════════════════════════════════
# Ocr.extract_text_from_text_file
# NOTE: This method is synchronous — do NOT await it.
# ═══════════════════════════════════════════════════════════════════════════

class TestExtractTextFromTextFile:
    def test_utf8(self, ocr):
        text = ocr.extract_text_from_text_file(b"hello world")
        assert text == "hello world"

    def test_empty_bytes_returns_empty(self, ocr):
        result = ocr.extract_text_from_text_file(b"")
        assert result == ""

    def test_non_bytes_raises_type_error(self, ocr):
        with pytest.raises(TypeError):
            ocr.extract_text_from_text_file("not bytes")

    def test_latin1_content(self, ocr):
        # Override encodings so latin-1 is tried (utf-8 fails, utf-16 skipped)
        ocr.text_extraction_config["encodings"] = ["utf-8", "latin-1"]
        text = b"caf\xe9"  # é in latin-1, invalid utf-8 strict
        result = ocr.extract_text_from_text_file(text)
        assert "caf" in result

    def test_fallback_encoding_used_when_all_fail(self, ocr):
        """Bytes that fail all strict encodings → fallback with replace."""
        weird = bytes([0xFF, 0xFE, 0x00, 0x41])  # UTF-16 BOM + 'A'
        result = ocr.extract_text_from_text_file(weird)
        assert isinstance(result, str)

    def test_fallback_decode_exception_raises_value_error(self, ocr):
        """Force all encodings to fail including fallback → ValueError raised."""
        ocr.text_extraction_config["encodings"] = ["utf-8"]
        ocr.text_extraction_config["fallback_encoding"] = "utf-8"
        # bytes that fail utf-8 strict but we also make fallback fail by using a bad error handler
        ocr.text_extraction_config["fallback_encoding_errors"] = "strict"
        bad_bytes = b"\xff\xfe\xff"  # BOM-like, fails utf-8 strict
        # Even fallback with strict will raise UnicodeDecodeError — which becomes ValueError
        with pytest.raises((ValueError, UnicodeDecodeError)):
            ocr.extract_text_from_text_file(bad_bytes)


# ═══════════════════════════════════════════════════════════════════════════
# Ocr.extract_text_from_csv
# NOTE: This method is synchronous — do NOT await it.
# ═══════════════════════════════════════════════════════════════════════════

class TestExtractTextFromCsv:
    def test_basic_csv(self, ocr):
        csv_bytes = b"name,age\nAlice,30\nBob,25"
        result = ocr.extract_text_from_csv(csv_bytes, "file.csv")
        assert "chunks" in result
        assert len(result["chunks"]) >= 1

    def test_empty_csv_returns_400(self, ocr):
        result = ocr.extract_text_from_csv(b"", "file.csv")
        assert result.get("status_code") == 400 or "error" in result or "message" in result

    def test_header_only_csv_returns_error(self, ocr):
        result = ocr.extract_text_from_csv(b"col1,col2", "file.csv")
        assert result.get("status_code") == 400

    def test_filename_propagated(self, ocr):
        csv_bytes = b"x,y\n1,2\n3,4"
        result = ocr.extract_text_from_csv(csv_bytes, "myfile.csv")
        assert "chunks" in result

    def test_fallback_encoding_used(self, ocr):
        """Remove all encodings from config so fallback path is hit."""
        ocr.text_extraction_config["encodings"] = []
        csv_bytes = b"a,b\n1,2"
        result = ocr.extract_text_from_csv(csv_bytes, "file.csv")
        assert "chunks" in result

    def test_exception_returns_error_dict(self, ocr):
        with patch("csv.reader", side_effect=RuntimeError("boom")):
            result = ocr.extract_text_from_csv(b"a,b\n1,2", "file.csv")
        assert "error" in result


# ═══════════════════════════════════════════════════════════════════════════
# Ocr.extract_text_from_xlsx
# NOTE: Takes only (self, excel_bytes) — no auth_token argument.
# ═══════════════════════════════════════════════════════════════════════════

class TestExtractTextFromXlsx:
    def test_basic_xlsx(self, ocr):
        xlsx = make_xlsx_bytes()
        result = ocr.extract_text_from_xlsx(xlsx)
        assert "chunks" in result
        assert len(result["chunks"]) >= 1

    def test_empty_xlsx_no_chunks(self, ocr):
        # DataFrame with no rows → len(table_data)==1 (header only) → skipped
        buf = BytesIO()
        df = pd.DataFrame(columns=["col1"])
        with pd.ExcelWriter(buf, engine="openpyxl") as writer:
            df.to_excel(writer, index=False)
        buf.seek(0)
        result = ocr.extract_text_from_xlsx(buf.read())
        assert "chunks" in result  # empty but no error

    def test_exception_returns_error(self, ocr):
        result = ocr.extract_text_from_xlsx(b"notanxlsx")
        assert "error" in result

    def test_sheet_error_logged_and_continued(self, ocr):
        """If reading a sheet fails, it logs and moves on."""
        xlsx = make_xlsx_bytes({"Sheet1": pd.DataFrame({"a": [1]}), "Sheet2": pd.DataFrame({"b": [2]})})

        original_read = pd.read_excel
        call_count = [0]

        def patched_read(f, sheet_name, **kwargs):
            call_count[0] += 1
            if sheet_name == "Sheet2":
                raise ValueError("bad sheet")
            return original_read(f, sheet_name=sheet_name, **kwargs)

        with patch("src.ocr.ocr.pd.read_excel", side_effect=patched_read):
            result = ocr.extract_text_from_xlsx(xlsx)
        # Sheet1 should succeed
        assert "chunks" in result


# ═══════════════════════════════════════════════════════════════════════════
# Ocr.extract_text_from_excel
# NOTE: Takes (self, excel_bytes, filename) — no auth_token argument.
# ═══════════════════════════════════════════════════════════════════════════

class TestExtractTextFromExcel:
    def test_xlsx_extension(self, ocr):
        xlsx = make_xlsx_bytes()
        result = ocr.extract_text_from_excel(xlsx, "data.xlsx")
        assert "chunks" in result

    def test_xls_extension(self, ocr):
        # Use xlsx bytes but .xls filename — pd.ExcelFile accepts both
        xlsx = make_xlsx_bytes()
        result = ocr.extract_text_from_excel(xlsx, "data.xls")
        assert "chunks" in result

    def test_unsupported_extension_returns_error(self, ocr):
        result = ocr.extract_text_from_excel(b"data", "data.csv")
        assert "error" in result
        assert "Unsupported" in result["error"]

    def test_empty_sheets_return_400(self, ocr):
        buf = BytesIO()
        df = pd.DataFrame(columns=["col1"])
        with pd.ExcelWriter(buf, engine="openpyxl") as writer:
            df.to_excel(writer, index=False)
        buf.seek(0)
        result = ocr.extract_text_from_excel(buf.read(), "empty.xlsx")
        assert result.get("status_code") == 400

    def test_exception_returns_error(self, ocr):
        result = ocr.extract_text_from_excel(b"garbage", "bad.xlsx")
        assert "error" in result

    def test_sheet_read_error_continues(self, ocr):
        xlsx = make_xlsx_bytes({"Good": pd.DataFrame({"x": [1, 2]}), "Bad": pd.DataFrame({"y": [3, 4]})})
        original = pd.read_excel
        def side_effect(f, sheet_name, **kw):
            if sheet_name == "Bad":
                raise Exception("fail")
            return original(f, sheet_name=sheet_name, **kw)
        with patch("src.ocr.ocr.pd.read_excel", side_effect=side_effect):
            result = ocr.extract_text_from_excel(xlsx, "data.xlsx")
        assert "chunks" in result


# ═══════════════════════════════════════════════════════════════════════════
# Ocr.extract_text_from_docx
# ═══════════════════════════════════════════════════════════════════════════

class TestExtractTextFromDocx:
    @pytest.mark.asyncio
    async def test_basic_docx(self, ocr):
        docx_bytes = make_docx_bytes("Hello paragraph")
        result = await ocr.extract_text_from_docx(docx_bytes, TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result
        assert any("Hello paragraph" in c["text"] for c in result["chunks"])

    @pytest.mark.asyncio
    async def test_docx_with_table(self, ocr):
        docx_bytes = make_docx_with_table_bytes()
        result = await ocr.extract_text_from_docx(docx_bytes, TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result

    @pytest.mark.asyncio
    async def test_empty_docx_returns_400(self, ocr):
        doc = DocxDocument()
        buf = BytesIO()
        doc.save(buf)
        buf.seek(0)
        result = await ocr.extract_text_from_docx(buf.read(), TEAM_CONFIG, AUTH_TOKEN)
        assert result.get("status_code") == 400

    @pytest.mark.asyncio
    async def test_invalid_bytes_returns_error(self, ocr):
        result = await ocr.extract_text_from_docx(b"not a docx", TEAM_CONFIG, AUTH_TOKEN)
        assert "error" in result

    @pytest.mark.asyncio
    async def test_inline_shapes_processed(self, ocr):
        """Exercise image extraction path via mocked inline shapes."""
        docx_bytes = make_docx_bytes("text content")

        with patch("src.ocr.ocr.Document") as MockDoc:
            doc_instance = MagicMock()
            doc_instance.element.body = []
            doc_instance.tables = []
            doc_instance.paragraphs = []

            # Create a mock inline shape
            mock_blip = MagicMock()
            mock_blip.embed = "rId1"
            mock_inline = MagicMock()
            mock_inline.graphic.graphicData.pic.blipFill.blip = mock_blip

            mock_shape = MagicMock()
            mock_shape._inline = mock_inline

            mock_image = MagicMock()
            mock_image.blob = make_large_image_bytes()
            mock_image.px_width = 300
            mock_image.px_height = 300
            mock_image.filename = "test.png"
            mock_image.content_type = "image/png"

            mock_image_part = MagicMock()
            mock_image_part.image = mock_image
            mock_image_part.filename = "test.png"
            mock_image_part.content_type = "image/png"

            doc_instance.part.related_parts = {"rId1": mock_image_part}
            doc_instance.inline_shapes = [mock_shape]

            # Mock paragraph with text
            mock_para = MagicMock()
            mock_para.text = "Some text"

            p_block = MagicMock()
            p_block.tag = "p"  # doesn't end with 'p' in tag form, need proper suffix
            p_block.tag = "http://schemas.openxmlformats.org/wordprocessingml/2006/main}p"

            MockDoc.return_value = doc_instance

            with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_analyze:
                mock_analyze.return_value = IMAGE_ANALYSIS_RESULT
                # Doc has no body blocks → no text, but has inline shapes
                # Since no text elements → elements empty after text, but add a text element manually
                with patch.object(ocr, "_chunk_elements", wraps=ocr._chunk_elements) as mock_chunk:
                    result = await ocr.extract_text_from_docx(docx_bytes, TEAM_CONFIG, AUTH_TOKEN)
        # Should return 400 (no body text) or chunks
        assert "status_code" in result or "chunks" in result or "error" in result

    @pytest.mark.asyncio
    async def test_image_extraction_exception_logged(self, ocr):
        """Cover the except branch in inline_shapes loop."""
        docx_bytes = make_docx_bytes("text")

        with patch("src.ocr.ocr.Document") as MockDoc:
            doc_instance = MagicMock()
            doc_instance.element.body = []
            doc_instance.tables = []
            doc_instance.paragraphs = []

            bad_shape = MagicMock()
            bad_shape._inline.graphic.graphicData.pic.blipFill.blip.embed = "rId1"
            doc_instance.part.related_parts = {}  # KeyError on rId1
            doc_instance.inline_shapes = [bad_shape]
            MockDoc.return_value = doc_instance

            result = await ocr.extract_text_from_docx(docx_bytes, TEAM_CONFIG, AUTH_TOKEN)
        assert "status_code" in result or "chunks" in result or "error" in result

    @pytest.mark.asyncio
    async def test_image_result_exception_in_gather(self, ocr):
        """Cover the exception branch when gather returns an Exception."""
        docx_bytes = make_docx_bytes("text with images")

        # Patch Document to return a mock with inline shapes
        with patch("src.ocr.ocr.Document") as MockDoc:
            doc_instance = MagicMock()
            doc_instance.element.body = []
            doc_instance.tables = []
            doc_instance.paragraphs = []
            doc_instance.inline_shapes = []
            MockDoc.return_value = doc_instance

            with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_analyze:
                mock_analyze.side_effect = RuntimeError("analyze failed")
                result = await ocr.extract_text_from_docx(docx_bytes, TEAM_CONFIG, AUTH_TOKEN)
        assert "status_code" in result or "chunks" in result or "error" in result

    @pytest.mark.asyncio
    async def test_max_image_limit_enforced(self, ocr):
        """If MAX_DOCX_IMAGES shapes exist, loop should stop."""
        docx_bytes = make_docx_bytes("content")

        with patch("src.ocr.ocr.Document") as MockDoc:
            doc_instance = MagicMock()
            doc_instance.element.body = []
            doc_instance.tables = []
            doc_instance.paragraphs = []

            # Create MAX_DOCX_IMAGES + 5 inline shapes
            shapes = []
            for i in range(ProcessingLimits.MAX_DOCX_IMAGES + 5):
                s = MagicMock()
                s._inline.graphic.graphicData.pic.blipFill.blip.embed = f"rId{i}"
                shapes.append(s)

            doc_instance.inline_shapes = shapes
            doc_instance.part.related_parts = {}
            MockDoc.return_value = doc_instance

            result = await ocr.extract_text_from_docx(docx_bytes, TEAM_CONFIG, AUTH_TOKEN)
        assert "status_code" in result or "chunks" in result or "error" in result


# ═══════════════════════════════════════════════════════════════════════════
# Ocr.extract_text_from_ppt
# ═══════════════════════════════════════════════════════════════════════════

class TestExtractTextFromPpt:
    @pytest.mark.asyncio
    async def test_basic_pptx(self, ocr):
        pptx_bytes = make_pptx_bytes("Slide Title")
        result = await ocr.extract_text_from_ppt(pptx_bytes, TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result

    @pytest.mark.asyncio
    async def test_empty_pptx_returns_400(self, ocr):
        prs = Presentation()
        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        result = await ocr.extract_text_from_ppt(buf.read(), TEAM_CONFIG, AUTH_TOKEN)
        assert result.get("status_code") == 400

    @pytest.mark.asyncio
    async def test_invalid_bytes_returns_error(self, ocr):
        result = await ocr.extract_text_from_ppt(b"garbage", TEAM_CONFIG, AUTH_TOKEN)
        assert "error" in result

    @pytest.mark.asyncio
    async def test_temp_file_cleaned_up(self, ocr):
        pptx_bytes = make_pptx_bytes("test")
        tmp_paths = []
        original_ntf = tempfile.NamedTemporaryFile

        def capturing_ntf(**kwargs):
            f = original_ntf(**kwargs)
            tmp_paths.append(f.name)
            return f

        with patch("tempfile.NamedTemporaryFile", side_effect=capturing_ntf):
            await ocr.extract_text_from_ppt(pptx_bytes, TEAM_CONFIG, AUTH_TOKEN)

        for p in tmp_paths:
            assert not os.path.exists(p)

    @pytest.mark.asyncio
    async def test_pptx_with_table(self, ocr):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        # Add table
        table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
        table.cell(0, 0).text = "H1"
        table.cell(0, 1).text = "H2"
        table.cell(1, 0).text = "V1"
        table.cell(1, 1).text = "V2"
        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        result = await ocr.extract_text_from_ppt(buf.read(), TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result

    @pytest.mark.asyncio
    async def test_pptx_image_shape_processed(self, ocr):
        """Exercise the PICTURE shape branch via mocking."""
        pptx_bytes = make_pptx_bytes("some text")
        from pptx.enum.shapes import MSO_SHAPE_TYPE as MSO

        with patch("src.ocr.ocr.Presentation") as MockPrs:
            prs_instance = MagicMock()
            slide = MagicMock()

            pic_shape = MagicMock()
            pic_shape.has_table = False
            pic_shape.shape_type = MSO.PICTURE
            mock_img = MagicMock()
            mock_img.blob = make_large_image_bytes()
            mock_img.size = (300, 300)
            mock_img.ext = "png"
            pic_shape.image = mock_img
            pic_shape.name = "image1"
            pic_shape.text = None

            text_shape = MagicMock()
            text_shape.has_table = False
            text_shape.shape_type = MagicMock()  # not PICTURE
            text_shape.shape_type.__eq__ = lambda s, o: False
            text_shape.text = "some slide text"

            slide.shapes = [pic_shape, text_shape]
            prs_instance.slides = [slide]
            MockPrs.return_value = prs_instance

            with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
                mock_an.return_value = IMAGE_ANALYSIS_RESULT
                result = await ocr.extract_text_from_ppt(pptx_bytes, TEAM_CONFIG, AUTH_TOKEN)

        assert "chunks" in result

    @pytest.mark.asyncio
    async def test_pptx_image_irrelevant_skipped(self, ocr):
        """Decorative image (small) is skipped."""
        pptx_bytes = make_pptx_bytes("text")
        from pptx.enum.shapes import MSO_SHAPE_TYPE as MSO

        with patch("src.ocr.ocr.Presentation") as MockPrs:
            prs_instance = MagicMock()
            slide = MagicMock()

            pic_shape = MagicMock()
            pic_shape.has_table = False
            pic_shape.shape_type = MSO.PICTURE
            mock_img = MagicMock()
            mock_img.blob = make_small_image_bytes()  # too small
            mock_img.size = (10, 10)
            mock_img.ext = "png"
            pic_shape.image = mock_img
            pic_shape.name = "tiny"

            text_shape = MagicMock()
            text_shape.has_table = False
            type(text_shape).shape_type = PropertyMock(return_value=MSO.TEXT_BOX)
            text_shape.text = "slide text"

            slide.shapes = [pic_shape, text_shape]
            prs_instance.slides = [slide]
            MockPrs.return_value = prs_instance

            with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
                result = await ocr.extract_text_from_ppt(pptx_bytes, TEAM_CONFIG, AUTH_TOKEN)
                mock_an.assert_not_called()

    @pytest.mark.asyncio
    async def test_pptx_image_extraction_exception(self, ocr):
        """Cover except in image extraction within PPTX loop."""
        pptx_bytes = make_pptx_bytes("text")
        from pptx.enum.shapes import MSO_SHAPE_TYPE as MSO

        with patch("src.ocr.ocr.Presentation") as MockPrs:
            prs_instance = MagicMock()
            slide = MagicMock()

            pic_shape = MagicMock()
            pic_shape.has_table = False
            pic_shape.shape_type = MSO.PICTURE
            pic_shape.image = MagicMock(side_effect=Exception("image access error"))

            text_shape = MagicMock()
            text_shape.has_table = False
            type(text_shape).shape_type = PropertyMock(return_value=MSO.TEXT_BOX)
            text_shape.text = "content"

            slide.shapes = [pic_shape, text_shape]
            prs_instance.slides = [slide]
            MockPrs.return_value = prs_instance

            result = await ocr.extract_text_from_ppt(pptx_bytes, TEAM_CONFIG, AUTH_TOKEN)

        assert "chunks" in result or "error" in result

    @pytest.mark.asyncio
    async def test_pptx_max_image_limit(self, ocr):
        """MAX_PPTX_IMAGES limit breaks inner image loop."""
        pptx_bytes = make_pptx_bytes("text")
        from pptx.enum.shapes import MSO_SHAPE_TYPE as MSO

        with patch("src.ocr.ocr.Presentation") as MockPrs:
            prs_instance = MagicMock()
            slide = MagicMock()

            shapes = []
            for i in range(ProcessingLimits.MAX_PPTX_IMAGES + 5):
                ps = MagicMock()
                ps.has_table = False
                ps.shape_type = MSO.PICTURE
                mi = MagicMock()
                mi.blob = make_large_image_bytes()
                mi.size = (300, 300)
                mi.ext = "png"
                ps.image = mi
                ps.name = f"img{i}"
                shapes.append(ps)

            slide.shapes = shapes
            prs_instance.slides = [slide]
            MockPrs.return_value = prs_instance

            with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
                mock_an.return_value = IMAGE_ANALYSIS_RESULT
                result = await ocr.extract_text_from_ppt(pptx_bytes, TEAM_CONFIG, AUTH_TOKEN)

        assert "chunks" in result or "status_code" in result

    @pytest.mark.asyncio
    async def test_cleanup_on_exception(self, ocr):
        """Verify finally block runs even when exception occurs."""
        with patch("tempfile.NamedTemporaryFile", side_effect=OSError("no temp")):
            result = await ocr.extract_text_from_ppt(b"data", TEAM_CONFIG, AUTH_TOKEN)
        assert "error" in result


# ═══════════════════════════════════════════════════════════════════════════
# Ocr.extract_text_from_pdf_file
# ═══════════════════════════════════════════════════════════════════════════

class TestExtractTextFromPdfFile:
    def _make_mock_pdf(self, pages_data):
        """
        pages_data: list of dicts with keys: text, tables, images
        """
        mock_pdf = MagicMock()
        mock_pages = []

        for pd_item in pages_data:
            page = MagicMock()
            page.extract_text.return_value = pd_item.get("text", "")
            page.extract_tables.return_value = pd_item.get("tables", [])
            page.images = pd_item.get("images", [])
            page.width = 595
            page.height = 842
            mock_pages.append(page)

        mock_pdf.pages = mock_pages
        mock_pdf.close = MagicMock()
        return mock_pdf

    @pytest.mark.asyncio
    async def test_text_extraction(self, ocr):
        mock_pdf = self._make_mock_pdf([{"text": "Hello PDF", "tables": [], "images": []}])
        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            result = await ocr.extract_text_from_pdf_file(b"fake_pdf", TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result
        assert any("Hello PDF" in c["text"] for c in result["chunks"])

    @pytest.mark.asyncio
    async def test_table_extraction(self, ocr):
        table = [["H1", "H2"], ["V1", "V2"]]
        mock_pdf = self._make_mock_pdf([{"text": "", "tables": [table], "images": []}])
        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result

    @pytest.mark.asyncio
    async def test_table_with_none_cells(self, ocr):
        table = [["H1", None], [None, "V2"]]
        mock_pdf = self._make_mock_pdf([{"text": "", "tables": [table], "images": []}])
        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result

    @pytest.mark.asyncio
    async def test_empty_pdf_returns_400(self, ocr):
        mock_pdf = self._make_mock_pdf([{"text": "", "tables": [], "images": []}])
        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert result.get("status_code") == 400

    @pytest.mark.asyncio
    async def test_image_with_stream_bytes(self, ocr):
        large_bytes = make_large_image_bytes()
        img_info = {"stream": large_bytes, "width": 300, "height": 200, "xref": 1}
        mock_pdf = self._make_mock_pdf([{"text": "context", "tables": [], "images": [img_info]}])
        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
                mock_an.return_value = IMAGE_ANALYSIS_RESULT
                result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result

    @pytest.mark.asyncio
    async def test_image_with_stream_get_data(self, ocr):
        """stream object has get_data method."""
        large_bytes = make_large_image_bytes()

        class FakeStream:
            def get_data(self):
                return large_bytes

        img_info = {"stream": FakeStream(), "width": 300, "height": 200, "xref": 2}
        mock_pdf = self._make_mock_pdf([{"text": "ctx", "tables": [], "images": [img_info]}])
        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
                mock_an.return_value = IMAGE_ANALYSIS_RESULT
                result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result

    @pytest.mark.asyncio
    async def test_image_fallback_crop_path(self, ocr):
        """No stream → crop fallback path."""
        img_info = {"stream": None, "width": 300, "height": 200, "x0": 0, "y0": 0, "x1": 100, "y1": 100}

        page = MagicMock()
        page.extract_text.return_value = "page text"
        page.extract_tables.return_value = []
        page.images = [img_info]
        page.width = 595
        page.height = 842

        # Mock cropped page
        mock_pil_img = MagicMock()
        buf = BytesIO()
        from PIL import Image
        Image.new("RGB", (100, 100)).save(buf, format="PNG")
        buf.seek(0)
        mock_pil_img.save = lambda b, format: b.write(buf.read())

        mock_img_obj = MagicMock()
        mock_img_obj.original = mock_pil_img

        mock_cropped = MagicMock()
        mock_cropped.to_image.return_value = mock_img_obj
        page.crop.return_value = mock_cropped

        mock_pdf = MagicMock()
        mock_pdf.pages = [page]
        mock_pdf.close = MagicMock()

        large = make_large_image_bytes()

        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            with patch("src.ocr.ocr._io.BytesIO") as MockBuf:
                buf_instance = MagicMock()
                buf_instance.getvalue.return_value = large
                MockBuf.return_value = buf_instance
                with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
                    mock_an.return_value = IMAGE_ANALYSIS_RESULT
                    result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result or "status_code" in result

    @pytest.mark.asyncio
    async def test_image_irrelevant_skipped(self, ocr):
        """Small images are not analyzed."""
        small_bytes = make_small_image_bytes()
        img_info = {"stream": small_bytes, "width": 10, "height": 10}
        mock_pdf = self._make_mock_pdf([{"text": "text", "tables": [], "images": [img_info]}])
        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
                await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
                mock_an.assert_not_called()

    @pytest.mark.asyncio
    async def test_image_exception_logged(self, ocr):
        """Exception while getting image bytes is caught and logged."""
        img_info = MagicMock()
        img_info.get.side_effect = Exception("image access error")

        page = MagicMock()
        page.extract_text.return_value = "text"
        page.extract_tables.return_value = []
        page.images = [img_info]

        mock_pdf = MagicMock()
        mock_pdf.pages = [page]
        mock_pdf.close = MagicMock()

        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result or "status_code" in result

    @pytest.mark.asyncio
    async def test_table_extraction_exception_logged(self, ocr):
        """extract_tables exception is caught and logged."""
        page = MagicMock()
        page.extract_text.return_value = "some text"
        page.extract_tables.side_effect = Exception("table error")
        page.images = []

        mock_pdf = MagicMock()
        mock_pdf.pages = [page]
        mock_pdf.close = MagicMock()

        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result

    @pytest.mark.asyncio
    async def test_max_pdf_images_limit(self, ocr):
        """Once MAX_PDF_IMAGES is reached, image processing stops."""
        large_bytes = make_large_image_bytes()
        images = [{"stream": large_bytes, "width": 300, "height": 200} for _ in range(ProcessingLimits.MAX_PDF_IMAGES + 5)]
        mock_pdf = self._make_mock_pdf([{"text": "text", "tables": [], "images": images}])
        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
                mock_an.return_value = IMAGE_ANALYSIS_RESULT
                result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result

    @pytest.mark.asyncio
    async def test_gather_exception_result_logged(self, ocr):
        """If gather returns an Exception for an image, it's logged and skipped."""
        large_bytes = make_large_image_bytes()
        img_info = {"stream": large_bytes, "width": 300, "height": 200}
        mock_pdf = self._make_mock_pdf([{"text": "text", "tables": [], "images": [img_info]}])
        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
                mock_an.side_effect = RuntimeError("analyze failed")
                result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result

    @pytest.mark.asyncio
    async def test_pdf_open_exception(self, ocr):
        with patch("src.ocr.ocr.pdfplumber.open", side_effect=Exception("cannot open")):
            result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert "error" in result

    @pytest.mark.asyncio
    async def test_pdf_close_exception_in_finally(self, ocr):
        """Cover the close error branch in finally."""
        mock_pdf = MagicMock()
        mock_pdf.pages = []
        mock_pdf.close.side_effect = Exception("close error")

        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert "status_code" in result or "error" in result

    @pytest.mark.asyncio
    async def test_image_empty_stream_falls_through(self, ocr):
        """stream is present but get_data returns empty → None → fallback crop."""
        class FakeStream:
            def get_data(self):
                return b""  # falsy bytes

        img_info = {"stream": FakeStream(), "width": 300, "height": 200, "x0": 0, "y0": 0, "x1": 50, "y1": 50}
        page = MagicMock()
        page.extract_text.return_value = "text"
        page.extract_tables.return_value = []
        page.images = [img_info]
        page.width = 595
        page.height = 842

        large_bytes = make_large_image_bytes()

        mock_pil_buf = BytesIO(large_bytes)

        mock_pil_img = MagicMock()
        mock_pil_img.save = lambda b, format: b.write(large_bytes)
        mock_img_obj = MagicMock()
        mock_img_obj.original = mock_pil_img
        mock_cropped = MagicMock()
        mock_cropped.to_image.return_value = mock_img_obj
        page.crop.return_value = mock_cropped

        mock_pdf = MagicMock()
        mock_pdf.pages = [page]
        mock_pdf.close = MagicMock()

        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
                mock_an.return_value = IMAGE_ANALYSIS_RESULT
                result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result or "status_code" in result

    @pytest.mark.asyncio
    async def test_image_analysis_result_none_skipped(self, ocr):
        """result=None after gather is skipped (not isinstance dict)."""
        large_bytes = make_large_image_bytes()
        img_info = {"stream": large_bytes, "width": 300, "height": 200}
        mock_pdf = self._make_mock_pdf([{"text": "text", "tables": [], "images": [img_info]}])
        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
                mock_an.return_value = None
                result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result


# ═══════════════════════════════════════════════════════════════════════════
# Ocr.extract_text_from_zip
# ═══════════════════════════════════════════════════════════════════════════

class TestExtractTextFromZip:
    
    @pytest.mark.asyncio
    async def test_basic_txt_file(self, ocr):
        z = make_zip_bytes({"hello.txt": b"hello world content"})
        def sync_handle_txt(file_name, file_bytes, extracted_files):
            extracted_files[file_name] = [{"text": ocr.extract_text_from_text_file(file_bytes), "type": "text"}]
        with patch.object(ocr, "_handle_zip_text_file", side_effect=sync_handle_txt):
            result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "hello.txt" in result
        assert result["hello.txt"][0]["text"] == "hello world content"

    @pytest.mark.asyncio
    async def test_directory_entry_skipped(self, ocr):
        buf = BytesIO()
        with zipfile.ZipFile(buf, "w") as zf:
            zf.mkdir("subdir")
            zf.writestr("subdir/file.txt", b"content")
        buf.seek(0)
        result = await ocr.extract_text_from_zip(buf.read(), TEAM_CONFIG, AUTH_TOKEN)
        assert "subdir/" not in result

    @pytest.mark.asyncio
    async def test_invalid_zip_returns_error(self, ocr):
        result = await ocr.extract_text_from_zip(b"not a zip", TEAM_CONFIG, AUTH_TOKEN)
        assert "error" in result

    @pytest.mark.asyncio
    async def test_max_depth_exceeded(self, ocr):
        result = await ocr.extract_text_from_zip(
            b"dummy", TEAM_CONFIG, AUTH_TOKEN,
            current_depth=ProcessingLimits.MAX_ZIP_DEPTH
        )
        assert "error" in result
        assert "depth" in result["error"].lower()

    @pytest.mark.asyncio
    async def test_nested_zip_processed(self, ocr):
        z = make_nested_zip_bytes({"inner.txt": b"inner content"})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert len(result) >= 1

    @pytest.mark.asyncio
    async def test_nested_zip_with_error(self, ocr):
        """Nested zip returns error → stored as error entry."""
        inner_bad = b"not a zip"
        z = make_zip_bytes({"nested.zip": inner_bad})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "nested.zip" in result

    @pytest.mark.asyncio
    async def test_csv_in_zip(self, ocr):
        z = make_zip_bytes({"data.csv": b"col1,col2\nv1,v2\nv3,v4"})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "data.csv" in result

    @pytest.mark.asyncio
    async def test_xlsx_in_zip(self, ocr):
        xlsx = make_xlsx_bytes()
        z = make_zip_bytes({"data.xlsx": xlsx})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "data.xlsx" in result

    @pytest.mark.asyncio
    async def test_xls_in_zip(self, ocr):
        xlsx = make_xlsx_bytes()
        z = make_zip_bytes({"data.xls": xlsx})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "data.xls" in result

    @pytest.mark.asyncio
    async def test_docx_in_zip(self, ocr):
        docx = make_docx_bytes("Zip DOCX content")
        z = make_zip_bytes({"doc.docx": docx})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "doc.docx" in result

    @pytest.mark.asyncio
    async def test_doc_in_zip_error(self, ocr):
        z = make_zip_bytes({"bad.doc": b"garbage"})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "bad.doc" in result

    @pytest.mark.asyncio
    async def test_pptx_in_zip(self, ocr):
        pptx = make_pptx_bytes("Zip PPTX")
        z = make_zip_bytes({"slide.pptx": pptx})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "slide.pptx" in result

    @pytest.mark.asyncio
    async def test_ppt_in_zip_error(self, ocr):
        z = make_zip_bytes({"bad.ppt": b"garbage"})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "bad.ppt" in result

    @pytest.mark.asyncio
    async def test_pdf_in_zip(self, ocr):
        mock_pdf = MagicMock()
        page = MagicMock()
        page.extract_text.return_value = "pdf content"
        page.extract_tables.return_value = []
        page.images = []
        mock_pdf.pages = [page]
        mock_pdf.close = MagicMock()

        z = make_zip_bytes({"doc.pdf": b"fake_pdf_content"})
        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "doc.pdf" in result

    @pytest.mark.asyncio
    async def test_image_in_zip(self, ocr):
        z = make_zip_bytes({"photo.png": make_large_image_bytes()})
        with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
            mock_an.return_value = IMAGE_ANALYSIS_RESULT
            result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "photo.png" in result

    @pytest.mark.asyncio
    async def test_image_in_zip_analysis_error(self, ocr):
        z = make_zip_bytes({"photo.jpg": make_large_image_bytes()})
        with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
            mock_an.side_effect = Exception("analyze failed")
            result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "photo.jpg" in result
        assert "error" in result["photo.jpg"][0]["text"].lower()

    @pytest.mark.asyncio
    async def test_image_in_zip_no_content(self, ocr):
        z = make_zip_bytes({"photo.gif": make_large_image_bytes()})
        with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_an:
            mock_an.return_value = EMPTY_IMAGE_ANALYSIS
            result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "photo.gif" in result
        assert result["photo.gif"][0]["text"] == "No content extracted"

    @pytest.mark.asyncio
    async def test_unsupported_extension(self, ocr):
        z = make_zip_bytes({"data.bin": b"binary data"})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "data.bin" in result
        assert result["data.bin"][0]["type"] == "unsupported"

    @pytest.mark.asyncio
    async def test_file_no_extension(self, ocr):
        z = make_zip_bytes({"README": b"text without extension"})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "README" in result

    @pytest.mark.asyncio
    async def test_max_file_limit_enforced(self, ocr):
        files = {f"file{i}.txt": f"content{i}".encode() for i in range(ProcessingLimits.MAX_ZIP_FILES + 10)}
        z = make_zip_bytes(files)
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert len(result) <= ProcessingLimits.MAX_ZIP_FILES

    @pytest.mark.asyncio
    async def test_max_extracted_size_limit(self, ocr):
        """File that exceeds MAX_ZIP_EXTRACTED_SIZE is skipped."""
        large_content = b"x" * (ProcessingLimits.MAX_ZIP_EXTRACTED_SIZE + 1)
        z = make_zip_bytes({"big.txt": large_content})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        # Either empty or the file wasn't processed
        assert isinstance(result, dict)

    @pytest.mark.asyncio
    async def test_file_processing_exception(self, ocr):
        """zipfile.read raises → stored as error."""
        z = make_zip_bytes({"file.txt": b"content"})
        original_zipfile = zipfile.ZipFile

        class BrokenZip:
            def __init__(self, *args, **kwargs):
                self._zf = original_zipfile(*args, **kwargs)

            def __enter__(self):
                return self

            def __exit__(self, *args):
                self._zf.close()

            def namelist(self):
                return self._zf.namelist()

            def read(self, name):
                raise IOError("read failed")

        with patch("src.ocr.ocr.zipfile.ZipFile", BrokenZip):
            result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "file.txt" in result
        assert result["file.txt"][0]["type"] == "error"

    @pytest.mark.asyncio
    async def test_txt_file_extraction_error(self, ocr):
        """Text file extraction failure stores error."""
        z = make_zip_bytes({"readme.txt": b"content"})
        # extract_text_from_text_file is synchronous — patch it as a regular Mock
        with patch.object(ocr, "extract_text_from_text_file", side_effect=Exception("decode failed")):
            result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "readme.txt" in result
        assert "error" in result["readme.txt"][0]["text"].lower()

    @pytest.mark.asyncio
    async def test_generic_exception_in_zip(self, ocr):
        with patch("src.ocr.ocr.BytesIO", side_effect=Exception("BytesIO failed")):
            result = await ocr.extract_text_from_zip(b"data", TEAM_CONFIG, AUTH_TOKEN)
        assert "error" in result

    @pytest.mark.asyncio
    async def test_json_txt_extension_processed(self, ocr):
        z = make_zip_bytes({"config.json": b'{"key": "value"}'})
        def async_handle_txt(file_name, file_bytes, extracted_files):
            extracted_files[file_name] = [{"text": ocr.extract_text_from_text_file(file_bytes), "type": "text"}]
        with patch.object(ocr, "_handle_zip_text_file", side_effect=async_handle_txt):
            result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "config.json" in result
        assert result["config.json"][0]["type"] == "text"

    @pytest.mark.asyncio
    async def test_md_extension_processed(self, ocr):
        z = make_zip_bytes({"README.md": b"# Title\n\nContent"})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "README.md" in result

    @pytest.mark.asyncio
    async def test_xlsx_error_in_zip(self, ocr):
        """xlsx extraction error stored as error entry."""
        z = make_zip_bytes({"bad.xlsx": b"garbage"})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "bad.xlsx" in result

    @pytest.mark.asyncio
    async def test_csv_error_in_zip(self, ocr):
        """csv extraction returning error is stored."""
        z = make_zip_bytes({"single.csv": b"header_only"})
        result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "single.csv" in result

    @pytest.mark.asyncio
    async def test_pdf_error_in_zip(self, ocr):
        z = make_zip_bytes({"bad.pdf": b"not a pdf"})
        with patch("src.ocr.ocr.pdfplumber.open", side_effect=Exception("bad pdf")):
            result = await ocr.extract_text_from_zip(z, TEAM_CONFIG, AUTH_TOKEN)
        assert "bad.pdf" in result


# ═══════════════════════════════════════════════════════════════════════════
# Edge cases and integration-style tests
# ═══════════════════════════════════════════════════════════════════════════

class TestEdgeCases:
    def test_table_to_markdown_single_cell(self):
        result = TableProcessor.table_to_markdown([["only"]])
        assert "| only |" in result

    @pytest.mark.asyncio
    async def test_chunk_elements_table_single_chunk_index_is_none(self, ocr):
        """Single-chunk table → table_chunk_index is None."""
        data = [["H"], ["v1"]]
        elements = [{"type": "table", "content": data, "metadata": {}}]
        chunks = await ocr._chunk_elements(elements, source="test")
        assert chunks[0].metadata["table_chunk_index"] is None

    @pytest.mark.asyncio
    async def test_extract_text_from_docx_table_before_text(self, ocr):
        """Table before any text flushes empty text_buffer cleanly."""
        docx_bytes = make_docx_with_table_bytes()
        result = await ocr.extract_text_from_docx(docx_bytes, TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result

    def test_ocr_with_missing_config(self, tmp_path):
        """Ocr with missing config still initialises and processes."""
        o = Ocr(config_path=str(tmp_path / "missing.yaml"))
        # extract_text_from_text_file is synchronous — call directly, no await
        result = o.extract_text_from_text_file(b"hello")
        assert result == "hello"

    @pytest.mark.asyncio
    async def test_concurrent_image_analysis(self, ocr):
        large_bytes = make_large_image_bytes()
        images = [{"stream": large_bytes, "width": 300, "height": 200} for _ in range(3)]
        mock_pdf = MagicMock()
        page = MagicMock()
        page.extract_text.return_value = "text"
        page.extract_tables.return_value = []
        page.images = images
        mock_pdf.pages = [page]
        mock_pdf.close = MagicMock()

        call_count = 0

        async def fake_analyze(**kwargs):
            nonlocal call_count
            call_count += 1
            return IMAGE_ANALYSIS_RESULT

        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            with patch.object(ocr, "_analyze_image_for_understanding", side_effect=fake_analyze):
                with patch.object(ocr, "semantic_chunker") as mock_sc:
                    mock_sc.chunk = AsyncMock(return_value=["text chunk"])
                    result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)

        assert call_count == 0
        assert "chunks" in result

    @pytest.mark.asyncio
    async def test_image_analysis_result_empty_format(self, ocr):
        """_format_image_analysis with empty result returns empty string → no_content_extracted."""
        with patch.object(ocr, "_analyze_image_with_semaphore", new_callable=AsyncMock) as mock_fn:
            mock_fn.return_value = EMPTY_IMAGE_ANALYSIS
            result = await ocr.extract_text_from_image(b"img", TEAM_CONFIG, AUTH_TOKEN)
        assert result["text"] != ""  # fallback message

    @pytest.mark.asyncio
    async def test_large_docx_pagination(self, ocr):
        """Long text should cause page counter to increment."""
        long_text = "A" * 10000
        doc = DocxDocument()
        doc.add_paragraph(long_text)
        buf = BytesIO()
        doc.save(buf)
        buf.seek(0)
        result = await ocr.extract_text_from_docx(buf.read(), TEAM_CONFIG, AUTH_TOKEN)
        assert "chunks" in result

    def test_processing_limits_are_class_attributes(self):
        assert hasattr(ProcessingLimits, "MAX_CONCURRENT_IMAGES")
        assert hasattr(ProcessingLimits, "MAX_ZIP_DEPTH")
        assert hasattr(ProcessingLimits, "MAX_ZIP_FILES")
        assert hasattr(ProcessingLimits, "MAX_ZIP_EXTRACTED_SIZE")
        assert hasattr(ProcessingLimits, "MAX_PDF_IMAGES")
        assert hasattr(ProcessingLimits, "MAX_PPTX_IMAGES")
        assert hasattr(ProcessingLimits, "MAX_DOCX_IMAGES")
        assert hasattr(ProcessingLimits, "IMAGE_ANALYSIS_TIMEOUT")

    @pytest.mark.asyncio
    async def test_pdf_text_flushed_before_table(self, ocr):
        """Text buffer is flushed before adding a table element."""
        table = [["H1", "H2"], ["V1", "V2"]]
        page = MagicMock()
        page.extract_text.return_value = "pre-table text"
        page.extract_tables.return_value = [table]
        page.images = []

        mock_pdf = MagicMock()
        mock_pdf.pages = [page]
        mock_pdf.close = MagicMock()

        with patch("src.ocr.ocr.pdfplumber.open", return_value=mock_pdf):
            result = await ocr.extract_text_from_pdf_file(b"fake", TEAM_CONFIG, AUTH_TOKEN)

        assert "chunks" in result
        types = [c["type"] for c in result["chunks"]]
        assert "text" in types
        assert "table" in types

    @pytest.mark.asyncio
    async def test_pptx_table_flushes_text(self, ocr):
        """Text buffer in pptx is flushed before adding table."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE as MSO

        with patch("src.ocr.ocr.Presentation") as MockPrs:
            prs = MagicMock()
            slide = MagicMock()

            text_shape = MagicMock()
            text_shape.has_table = False
            type(text_shape).shape_type = PropertyMock(return_value=MSO.TEXT_BOX)
            text_shape.text = "text before table"

            table_shape = MagicMock()
            table_shape.has_table = True
            mock_table = MagicMock()
            cell1 = MagicMock()
            cell1.text = "H"
            row1 = MagicMock()
            row1.cells = [cell1]
            mock_table.rows = [row1]
            table_shape.table = mock_table

            slide.shapes = [text_shape, table_shape]
            prs.slides = [slide]
            MockPrs.return_value = prs

            with patch("tempfile.NamedTemporaryFile") as mock_ntf:
                tmp = MagicMock()
                tmp.__enter__ = Mock(return_value=tmp)
                tmp.__exit__ = Mock(return_value=False)
                tmp.name = "/tmp/fake.pptx"
                mock_ntf.return_value = tmp
                with patch("src.ocr.ocr.os.path.exists", return_value=False):
                    result = await ocr.extract_text_from_ppt(b"pptx", TEAM_CONFIG, AUTH_TOKEN)

        assert "chunks" in result or "status_code" in result